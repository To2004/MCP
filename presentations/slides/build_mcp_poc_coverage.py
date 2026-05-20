"""Generate 2026-05-19_mcp-poc-coverage.pptx.

Single slide, text only — no shapes, colors, or decorations.
Layout: title at top, then 3 side-by-side text columns (one per MCP).
Each column: MCP name, one-sentence description, 3 tools with brief notes.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "2026-05-19_mcp-poc-coverage-v2.pptx"

BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY  = RGBColor(0x44, 0x44, 0x44)

# ── Geometry ─────────────────────────────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

MARGIN  = Inches(0.50)
GAP     = Inches(0.30)
AVAIL_W = SW - 2 * MARGIN - 2 * GAP
COL_W   = AVAIL_W // 3
COL_X   = [
    MARGIN,
    MARGIN + COL_W + GAP,
    MARGIN + 2 * COL_W + 2 * GAP,
]

# ── Helper ────────────────────────────────────────────────────────────────────

def add_para(tf, text, size, bold=False, italic=False, color=BLACK, space_before=0):
    """Append one paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color


def col_textbox(slide, col_idx, top, height):
    """Return the text frame of a transparent text box in the given column."""
    tb = slide.shapes.add_textbox(COL_X[col_idx], top, COL_W, height)
    tb.text_frame.word_wrap = True
    # Remove the default empty first paragraph so callers control all content
    tb.text_frame.clear()
    return tb.text_frame


# ── Content ───────────────────────────────────────────────────────────────────

MCPS = [
    {
        "name": "Filesystem MCP",
        "desc": (
            "An MCP server that exposes a local directory tree, allowing "
            "an AI agent to read, write, and manipulate files on the host machine."
        ),
        "tools": [
            ("read_file",
             "Returns the full content of any file — can expose credentials, "
             "keys, or source code stored anywhere in the tree."),
            ("write_file",
             "Creates or overwrites a file with agent-supplied content — "
             "can inject backdoors or corrupt critical configs."),
            ("move_file",
             "Relocates a file to a new path — can hide sensitive files "
             "or stage data for later exfiltration."),
        ],
    },
    {
        "name": "SQLite MCP",
        "desc": (
            "An MCP server wrapping a SQLite database, enabling an AI agent "
            "to query and modify structured data through SQL operations."
        ),
        "tools": [
            ("read_query",
             "Executes SELECT statements — grants read access to any table, "
             "exposing PII, credentials, or business records at scale."),
            ("write_query",
             "Executes INSERT / UPDATE / DELETE — three blast radii under "
             "a single tool name; can corrupt or wipe stored data."),
            ("list_tables",
             "Returns all table names — full schema discovery with one call, "
             "the first step in any targeted data attack."),
        ],
    },
    {
        "name": "Slack MCP",
        "desc": (
            "An MCP server that connects an AI agent to a Slack workspace, "
            "enabling it to read messages, post content, and interact with channels."
        ),
        "tools": [
            ("post_message",
             "Sends messages as the authenticated user — enables phishing, "
             "misinformation, or data exfiltration via chat."),
            ("get_channel_history",
             "Retrieves full conversation history — exposes private discussions, "
             "credentials shared in chat, and strategic decisions."),
            ("list_channels",
             "Enumerates all visible channels — maps the org's communication "
             "structure for targeted follow-up attacks."),
        ],
    },
]


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Main title ────────────────────────────────────────────────────
    title_tb = slide.shapes.add_textbox(MARGIN, Inches(0.18), SW - 2 * MARGIN, Inches(0.65))
    tf = title_tb.text_frame
    tf.word_wrap = False
    tf.clear()
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = (
        "We started with 3 MCPs to cover a wide range of operations "
        "and misuse possibilities"
    )
    r.font.size  = Pt(22)
    r.font.bold  = True
    r.font.color.rgb = BLACK

    # ── Columns ───────────────────────────────────────────────────────
    col_top = Inches(1.00)
    col_h   = SH - col_top - Inches(0.20)

    for ci, mcp in enumerate(MCPS):
        tf = col_textbox(slide, ci, col_top, col_h)

        # MCP name
        add_para(tf, mcp["name"], 18, bold=True)

        # One-sentence description
        add_para(tf, mcp["desc"], 12, italic=True, color=GRAY, space_before=6)

        # Tools
        for tool_name, tool_desc in mcp["tools"]:
            add_para(tf, tool_name, 13, bold=True, space_before=14)
            add_para(tf, tool_desc, 11, color=GRAY, space_before=2)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
