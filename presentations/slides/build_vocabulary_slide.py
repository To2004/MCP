"""Generate 2026-05-19_mcp-tool-vocabulary.pptx.

One slide: "Three MCPs, Six Operations" — maps each MCP's proprietary
tool names to the six atomic operations (read/write/modify/delete/
execute/configure), showing that the mapping is invisible to callers
and reviewers.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).parent / "2026-05-19_mcp-tool-vocabulary.pptx"

# ── Palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x17, 0x37, 0x5E)
BLUE    = RGBColor(0x1F, 0x49, 0x7D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF2, 0xF5, 0xF9)
DGRAY   = RGBColor(0x44, 0x44, 0x44)
MGRAY   = RGBColor(0xBB, 0xBB, 0xBB)

OP_COLORS = {
    "READ":      RGBColor(0x19, 0x6F, 0x3D),   # green
    "WRITE":     RGBColor(0x1A, 0x5C, 0xAB),   # blue
    "MODIFY":    RGBColor(0xC2, 0x6A, 0x00),   # amber
    "DELETE":    RGBColor(0xBF, 0x21, 0x1A),   # red
    "EXECUTE":   RGBColor(0x7B, 0x00, 0x00),   # dark red
    "CONFIGURE": RGBColor(0x5B, 0x2C, 0x8A),   # purple
}

OP_BG = {
    "READ":      RGBColor(0xD5, 0xF0, 0xE0),
    "WRITE":     RGBColor(0xD6, 0xE8, 0xF8),
    "MODIFY":    RGBColor(0xFD, 0xEF, 0xD8),
    "DELETE":    RGBColor(0xFA, 0xD7, 0xD5),
    "EXECUTE":   RGBColor(0xF5, 0xC6, 0xC4),
    "CONFIGURE": RGBColor(0xEA, 0xD9, 0xF8),
}

# ── Tool catalog ─────────────────────────────────────────────────────
ROWS = [
    ("READ",
     "read_file\nlist_directory\nsearch_files",
     "read_query\nlist_tables\ndescribe_table",
     "get_channel_history\nget_thread_replies\nget_users"),
    ("WRITE",
     "write_file\ncreate_directory",
     "write_query (INSERT)",
     "post_message\nreply_to_thread"),
    ("MODIFY",
     "edit_file\nmove_file",
     "write_query (UPDATE)",
     "update_message\nadd_reaction"),
    ("DELETE",
     "delete_file",
     "write_query (DELETE)",
     "delete_message"),
    ("EXECUTE",
     "run_shell_command",
     "—",
     "—"),
    ("CONFIGURE",
     "create_directory",
     "create_table",
     "create_channel\ninvite_to_channel"),
]

HEADERS = ["", "Filesystem MCP", "SQLite MCP", "Slack MCP"]

# ── Slide geometry ────────────────────────────────────────────────────
SW, SH = Inches(13.333), Inches(7.5)   # 16:9

MARGIN_L  = Inches(0.45)
MARGIN_T  = Inches(0.20)
TITLE_H   = Inches(1.05)
FOOTER_H  = Inches(0.42)
TABLE_T   = MARGIN_T + TITLE_H + Inches(0.08)
TABLE_H   = SH - TABLE_T - FOOTER_H - Inches(0.15)
TABLE_W   = SW - MARGIN_L * 2

COL_WIDTHS = [Inches(1.60), Inches(3.49), Inches(3.49), Inches(3.85)]

HDR_H   = Inches(0.46)
DATA_ROW_H = (TABLE_H - HDR_H) / len(ROWS)


def _solid(shape, rgb: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _para(tf, text: str, size: int, bold: bool, color: RGBColor,
          align=PP_ALIGN.LEFT, italic: bool = False, mono: bool = False):
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if mono:
            run.font.name = "Courier New"


def cell_tf(table, row, col):
    return table.cell(row, col).text_frame


def build():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    blank_layout = prs.slide_layouts[6]   # completely blank
    slide = prs.slides.add_slide(blank_layout)
    shapes = slide.shapes

    # ── Background ──────────────────────────────────────────────────
    bg = shapes.add_shape(1, 0, 0, SW, SH)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    _solid(bg, LGRAY)
    bg.line.fill.background()

    # ── Title bar ──────────────────────────────────────────────────
    title_bar = shapes.add_shape(1, 0, 0, SW, TITLE_H)
    _solid(title_bar, NAVY)
    title_bar.line.fill.background()

    title_box = shapes.add_textbox(MARGIN_L, Inches(0.10), SW - MARGIN_L * 2, Inches(0.52))
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Three MCPs, Six Operations"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE

    sub_box = shapes.add_textbox(MARGIN_L, Inches(0.62), SW - MARGIN_L * 2, Inches(0.36))
    tf2 = sub_box.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = (
        "Every tool name resolves to one of six atomic operations — "
        "but the mapping is invisible to callers and reviewers."
    )
    r2.font.size = Pt(13)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor(0xC8, 0xD8, 0xEE)

    # ── Main table ─────────────────────────────────────────────────
    n_rows = len(ROWS) + 1   # +1 header
    n_cols = 4
    table = shapes.add_table(n_rows, n_cols, MARGIN_L, TABLE_T,
                             TABLE_W, TABLE_H).table

    # Set column widths
    for ci, w in enumerate(COL_WIDTHS):
        table.columns[ci].width = w

    # Set row heights
    table.rows[0].height = HDR_H
    for ri in range(1, n_rows):
        table.rows[ri].height = int(DATA_ROW_H)

    # Header row
    for ci, label in enumerate(HEADERS):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE

    # Data rows
    for ri, (op, fs_tools, sql_tools, slack_tools) in enumerate(ROWS):
        row_idx = ri + 1
        tools = [fs_tools, sql_tools, slack_tools]

        # Operation label cell (col 0)
        op_cell = table.cell(row_idx, 0)
        op_cell.fill.solid()
        op_cell.fill.fore_color.rgb = OP_BG[op]
        tf = op_cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = op
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = OP_COLORS[op]

        # Tool cells (cols 1-3)
        for ci, tool_text in enumerate(tools, start=1):
            tcell = table.cell(row_idx, ci)
            bg_rgb = OP_BG[op] if (ri % 2 == 0) else RGBColor(0xFF, 0xFF, 0xFF)
            tcell.fill.solid()
            tcell.fill.fore_color.rgb = bg_rgb
            tf = tcell.text_frame
            tf.word_wrap = True
            for li, line in enumerate(tool_text.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(1)
                r = p.add_run()
                r.text = line
                r.font.size = Pt(11)
                r.font.name = "Courier New" if line not in ("—",) else "Calibri"
                r.font.color.rgb = DGRAY if line != "—" else MGRAY

    # ── Footer callout ─────────────────────────────────────────────
    footer_top = SH - FOOTER_H - Inches(0.08)
    footer_box = shapes.add_shape(1, MARGIN_L, footer_top,
                                  SW - MARGIN_L * 2, FOOTER_H)
    _solid(footer_box, RGBColor(0xFF, 0xF3, 0xCD))
    footer_box.line.color.rgb = RGBColor(0xC8, 0x9F, 0x00)

    ft_text = shapes.add_textbox(MARGIN_L + Inches(0.12), footer_top + Inches(0.07),
                                 SW - MARGIN_L * 2 - Inches(0.24), FOOTER_H - Inches(0.10))
    tf3 = ft_text.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    r3a = p3.add_run()
    r3a.text = "Key insight: "
    r3a.font.size = Pt(11)
    r3a.font.bold = True
    r3a.font.color.rgb = RGBColor(0x7B, 0x4F, 0x00)
    r3b = p3.add_run()
    r3b.text = (
        "SQLite's write_query handles INSERT, UPDATE, and DELETE — "
        "three different blast radii under a single name. "
        "Benign users trigger critical operations by accident; "
        "reviewers cannot audit tool calls at scale without a semantic layer."
    )
    r3b.font.size = Pt(11)
    r3b.font.color.rgb = RGBColor(0x5A, 0x3A, 0x00)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
