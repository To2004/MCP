"""Smoke tests for the filesystem teaching deck generator."""

import pytest
from pptx.dml.color import RGBColor

from scripts.generate_filesystem_deck import generate_deck


@pytest.fixture(scope="session")
def deck():
    return generate_deck()


def test_deck_has_18_slides(deck):
    assert len(deck.slides) == 18


def test_slide_backgrounds_are_dark(deck):
    for i, slide in enumerate(deck.slides):
        fill = slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0x1E, 0x1E, 0x2E), f"Slide {i+1} wrong bg"


def _first_text(slide) -> str:
    """Return the text of the first text-containing shape on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                return t
    return ""


def test_slide_titles(deck):
    slides = deck.slides
    assert "MCP Filesystem Server" in _first_text(slides[0])
    assert "What is MCP" in _first_text(slides[1])
    assert "Request" in _first_text(slides[2])
    assert "Security Boundary" in _first_text(slides[3])
    assert "Read Tools" in _first_text(slides[4])
    assert "read_text_file" in _first_text(slides[5])
    assert "read_media_file" in _first_text(slides[6])
    assert "read_multiple_files" in _first_text(slides[7])
    assert "Write Tools" in _first_text(slides[8])
    assert "Navigation" in _first_text(slides[9])
    assert "search_files" in _first_text(slides[10])
    assert "Info" in _first_text(slides[11])
    assert "Matrix" in _first_text(slides[12])
    assert "Glob" in _first_text(slides[13])
    assert "Path" in _first_text(slides[14])
    assert "Overwrite" in _first_text(slides[15])
    assert "Partial" in _first_text(slides[16])
    assert "Don't Exist" in _first_text(slides[17])
