"""Build FS Research Presentation — Detective Story Arc.

Story: probe → observe output → change tactic → probe again → discover CVE.
Audience: thesis supervisor.

Usage: uv run python tests/testbed/build_fs_presentation.py
Output: Literature_review/litreturereview.pptx  (OVERWRITTEN)
DO NOT TOUCH: Literature_review/litreturereview_original_backup.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "Literature_review" / "litreturereview_original_backup.pptx"
OUTPUT = ROOT / "Literature_review" / "litreturereview_new.pptx"

# ---------------------------------------------------------------------------
# Lenovo colour palette
# ---------------------------------------------------------------------------
RED = RGBColor(0xE2, 0x23, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
MID_GRAY = RGBColor(0x33, 0x33, 0x33)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_GREEN = RGBColor(0x98, 0xFF, 0x98)
CODE_YELLOW = RGBColor(0xFF, 0xE0, 0x80)
CODE_BLUE = RGBColor(0x80, 0xC8, 0xFF)

# ---------------------------------------------------------------------------
# Layout indices (from template inspection)
# ---------------------------------------------------------------------------
LO_TITLE_BLACK = 2       # Title Slide_Black — ph idx 0 title, idx 1 subtitle
LO_TITLE_ONLY_BLACK = 9  # Title Only_Black  — ph idx 0 title
LO_CLOSING_BLACK = 35    # Closing Slide_Black — no placeholders


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _add_run(para, text, color=WHITE, bold=False, size_pt=16, italic=False):
    run = para.add_run()
    run.text = text
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.size = Pt(size_pt)
    run.font.italic = italic
    return run


def add_slide(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_title(slide, text, size_pt=26):
    ph = slide.placeholders[0]
    ph.text = text
    for para in ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = WHITE
            run.font.bold = True
            run.font.size = Pt(size_pt)


def add_textbox(slide, left, top, width, height,
                text="", size_pt=16, bold=False, color=WHITE,
                align="left", italic=False, word_wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    if align == "center":
        para.alignment = PP_ALIGN.CENTER
    elif align == "right":
        para.alignment = PP_ALIGN.RIGHT
    _add_run(para, text, color=color, bold=bold, size_pt=size_pt, italic=italic)
    return txb


def add_multiline_textbox(slide, left, top, width, height,
                          lines: list[dict], word_wrap=True):
    """lines = [{'text': str, 'color': RGBColor, 'bold': bool, 'size_pt': int,
                 'italic': bool, 'space_before': int}]"""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    first = True
    for spec in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        if spec.get("space_before"):
            para.space_before = Pt(spec["space_before"])
        _add_run(
            para,
            spec.get("text", ""),
            color=spec.get("color", WHITE),
            bold=spec.get("bold", False),
            size_pt=spec.get("size_pt", 16),
            italic=spec.get("italic", False),
        )
    return txb


def add_rect(slide, left, top, width, height, fill_color=RED):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def code_block(slide, left, top, width, height,
               lines: list[tuple[str, RGBColor]]):
    """Dark background box with coloured monospace-style code lines."""
    add_rect(slide, left, top, width, height, fill_color=CODE_BG)
    txb = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.12),
        Inches(width - 0.24), Inches(height - 0.24)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for text, color in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        _add_run(para, text, color=color, size_pt=13)
    return txb


def divider(slide, y, color=RED):
    add_rect(slide, 0.83, y, 11.67, 0.03, fill_color=color)


def bullet_rows(slide, start_y, rows: list[tuple[str, str, str]],
                col1=0.83, col2=1.5, col3=5.5,
                w2=3.8, w3=7.0,
                s1=16, s2=16, s3=15):
    """rows = [(bullet_char, label, detail)]"""
    y = start_y
    for b, label, detail in rows:
        add_textbox(slide, col1, y, 0.55, 0.45, b, size_pt=s1, bold=True, color=RED)
        add_textbox(slide, col2, y, w2, 0.45, label, size_pt=s2, bold=True, color=WHITE)
        if detail:
            add_textbox(slide, col3, y, w3, 0.55, detail, size_pt=s3, color=LIGHT_GRAY)
        y += 0.55
    return y


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    s = add_slide(prs, LO_TITLE_BLACK)
    ph = s.placeholders[0]
    ph.text = "Probing the MCP Filesystem Server"
    for para in ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = WHITE
            run.font.bold = True
            run.font.size = Pt(34)
    ph2 = s.placeholders[1]
    ph2.text = "A systematic study of input / output behaviour and the CVE found along the way"
    for para in ph2.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = LIGHT_GRAY
            run.font.size = Pt(16)


def slide_02_what_is_mcp(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "What Is an MCP Server?")

    add_textbox(s, 0.83, 1.25, 11.6, 0.5,
                "Model Context Protocol — a standard that lets an AI agent call tools "
                "hosted on a remote server over JSON-RPC.",
                size_pt=17, color=WHITE)

    rows = [
        ("Agent",  "sends a tool-call request",
         "\"please run read_file with this path\""),
        ("Server", "executes the tool and returns a result",
         "file content, directory listing, or an error"),
        ("Boundary", "the server decides what the agent is allowed to touch",
         "sandbox, permissions, path validation"),
    ]
    y = 2.05
    for actor, action, detail in rows:
        add_textbox(s, 0.83, y, 1.6, 0.42, actor, size_pt=15, bold=True, color=RED)
        add_textbox(s, 2.5, y, 3.8, 0.42, action, size_pt=15, color=WHITE)
        add_textbox(s, 6.4, y, 6.1, 0.55, detail, size_pt=14, color=LIGHT_GRAY, italic=True)
        y += 0.6

    divider(s, 3.95, color=DARK_GRAY)

    add_textbox(s, 0.83, 4.1, 11.6, 0.5,
                "In this study: the server exposes two tools — read_file and list_directory.",
                size_pt=16, bold=True, color=WHITE)
    add_textbox(s, 0.83, 4.7, 11.6, 0.5,
                "The server is the PROTECTED ASSET. I am probing it the way a threat source would.",
                size_pt=15, color=RED)


def slide_03_json_protocol(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "How MCP Talks — The JSON Protocol")

    add_textbox(s, 0.83, 1.25, 11.6, 0.4,
                "Every interaction is a JSON-RPC 2.0 message. Here is what a tool call looks like:",
                size_pt=16, color=LIGHT_GRAY)

    # Request block
    add_textbox(s, 0.83, 1.75, 3.0, 0.35, "REQUEST  (agent → server)",
                size_pt=13, bold=True, color=RED)
    code_block(s, 0.83, 2.15, 5.5, 2.55, [
        ('{',                                              WHITE),
        ('  "jsonrpc": "2.0",',                           CODE_YELLOW),
        ('  "method": "tools/call",',                     CODE_YELLOW),
        ('  "params": {',                                  WHITE),
        ('    "name": "read_file",',                       CODE_BLUE),
        ('    "arguments": {',                             WHITE),
        ('      "path": "/tmp/mcp_sandbox/readme.txt"',   CODE_GREEN),
        ('    }',                                          WHITE),
        ('  },',                                           WHITE),
        ('  "id": 1',                                      WHITE),
        ('}',                                              WHITE),
    ])

    # Response block
    add_textbox(s, 6.83, 1.75, 3.5, 0.35, "RESPONSE  (server → agent)",
                size_pt=13, bold=True, color=RED)
    code_block(s, 6.83, 2.15, 5.5, 2.55, [
        ('{',                                              WHITE),
        ('  "jsonrpc": "2.0",',                           CODE_YELLOW),
        ('  "result": {',                                  WHITE),
        ('    "content": [',                               WHITE),
        ('      {',                                        WHITE),
        ('        "type": "text",',                        CODE_BLUE),
        ('        "text": "Version 1.0 ..."',              CODE_GREEN),
        ('      }',                                        WHITE),
        ('    ],',                                         WHITE),
        ('    "isError": false',                           CODE_YELLOW),
        ('  }',                                            WHITE),
        ('}',                                             WHITE),
    ])

    add_textbox(s, 0.83, 4.85, 11.6, 0.45,
                "An error response is identical in shape — isError: true, "
                "text contains the error message. Same envelope, different payload.",
                size_pt=14, color=LIGHT_GRAY)


def slide_04_two_tools(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "The Two Tools: read_file + list_directory")

    # read_file
    add_rect(s, 0.83, 1.25, 0.07, 3.8, fill_color=RED)
    add_textbox(s, 1.05, 1.25, 5.2, 0.4, "read_file",
                size_pt=19, bold=True, color=WHITE)
    add_multiline_textbox(s, 1.05, 1.75, 5.2, 3.2, [
        {"text": "Input",  "bold": True, "color": RED,        "size_pt": 14},
        {"text": '  path: string  — absolute path to a file', "color": LIGHT_GRAY, "size_pt": 14},
        {"text": " ", "size_pt": 6},
        {"text": "Output (success)", "bold": True, "color": RED, "size_pt": 14},
        {"text": '  text content of the file',               "color": LIGHT_GRAY, "size_pt": 14},
        {"text": " ", "size_pt": 6},
        {"text": "Output (error)",   "bold": True, "color": RED, "size_pt": 14},
        {"text": '  ENOENT, Access denied, TypeError …',     "color": LIGHT_GRAY, "size_pt": 14},
        {"text": " ", "size_pt": 6},
        {"text": "What I probed", "bold": True, "color": WHITE, "size_pt": 14},
        {"text": '  path content, type, encoding, length',   "color": LIGHT_GRAY, "size_pt": 14},
    ])

    # list_directory
    add_rect(s, 6.83, 1.25, 0.07, 3.8, fill_color=RED)
    add_textbox(s, 7.05, 1.25, 5.2, 0.4, "list_directory",
                size_pt=19, bold=True, color=WHITE)
    add_multiline_textbox(s, 7.05, 1.75, 5.2, 3.2, [
        {"text": "Input",  "bold": True, "color": RED,        "size_pt": 14},
        {"text": '  path: string  — absolute path to a dir',  "color": LIGHT_GRAY, "size_pt": 14},
        {"text": " ", "size_pt": 6},
        {"text": "Output (success)", "bold": True, "color": RED, "size_pt": 14},
        {"text": '  [FILE] name\n  [DIR] name  (one per line)',  "color": LIGHT_GRAY, "size_pt": 14},
        {"text": " ", "size_pt": 6},
        {"text": "Output (error)",   "bold": True, "color": RED, "size_pt": 14},
        {"text": '  ENOENT, Access denied …',                 "color": LIGHT_GRAY, "size_pt": 14},
        {"text": " ", "size_pt": 6},
        {"text": "What I probed", "bold": True, "color": WHITE, "size_pt": 14},
        {"text": '  dir paths, empty dirs, symlinks',         "color": LIGHT_GRAY, "size_pt": 14},
    ])

    add_textbox(s, 0.83, 5.3, 11.6, 0.45,
                "Both tools share the same path-validation layer "
                "— so every finding applies to both.",
                size_pt=15, color=LIGHT_GRAY, italic=True)


def _scenario_cube(slide, x, y, num, name, desc, n_inputs,
                   highlight=False):
    """Draw one scenario cube at (x, y) — 5.55 wide × 0.98 tall."""
    bg = RGBColor(0x2A, 0x06, 0x06) if highlight else RGBColor(0x1E, 0x1E, 0x1E)
    stripe = RED if highlight else DARK_GRAY
    name_color = RED if highlight else WHITE

    add_rect(slide, x, y, 5.55, 0.98, fill_color=bg)
    add_rect(slide, x, y, 0.07, 0.98, fill_color=stripe)
    add_textbox(slide, x + 0.15, y + 0.07, 0.65, 0.38,
                num, size_pt=12, bold=True, color=DARK_GRAY)
    add_textbox(slide, x + 0.85, y + 0.07, 3.6, 0.38,
                name, size_pt=14, bold=True, color=name_color)
    add_textbox(slide, x + 0.85, y + 0.50, 3.6, 0.36,
                desc, size_pt=11, color=LIGHT_GRAY)
    add_textbox(slide, x + 4.5, y + 0.07, 0.98, 0.38,
                f"{n_inputs} inputs", size_pt=11,
                color=RED if highlight else DARK_GRAY, align="right")


def slide_05_methodology(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "Experimental Design: Systematic Input Space Exploration")

    # ── section labels ──────────────────────────────────────────────────────
    add_textbox(s, 0.83, 1.2, 5.55, 0.38,
                "PHASE I — CHARACTERISATION", size_pt=15, bold=True, color=WHITE)
    add_textbox(s, 0.83, 1.55, 5.55, 0.3,
                "Establish nominal server behaviour and response envelope",
                size_pt=11, color=LIGHT_GRAY, italic=True)

    add_textbox(s, 7.1, 1.2, 5.55, 0.38,
                "PHASE II — BOUNDARY PROBING", size_pt=15, bold=True, color=RED)
    add_textbox(s, 7.1, 1.55, 5.55, 0.3,
                "Exploit observations from Phase I to stress-test the access-control policy",
                size_pt=11, color=LIGHT_GRAY, italic=True)

    # vertical divider
    add_rect(s, 6.66, 1.15, 0.04, 6.05, fill_color=MID_GRAY)

    # ── characterisation cubes ───────────────────────────────────────────────
    baseline = [
        ("S1", "Baseline Characterisation",
         "Enumerate nominal read & list responses across valid path classes", 13),
        ("S2", "Error Response Analysis",
         "Map error message structure for missing, denied, and unreachable paths", 8),
        ("S3", "Path Resolution Behaviour",
         "Determine whether canonical form is derived before or after ACL evaluation", 8),
        ("S4", "Malformed Input Handling",
         "Assess server resilience to syntactically invalid path strings", 14),
    ]
    y = 2.0
    for num, name, desc, n in baseline:
        _scenario_cube(s, 0.83, y, num, name, desc, n, highlight=False)
        y += 1.08

    # ── boundary probing cubes ───────────────────────────────────────────────
    inference = [
        ("S5", "Near-Boundary Input Testing",
         "Submit near-valid paths to identify unhandled edge cases at the policy boundary",
         14, False),
        ("S6", "Semantic Boundary Violation",
         "Test whether a string-prefix ACL check can be circumvented without traversal",
         16, True),
        ("S7", "Type Coercion Analysis",
         "Supply non-string JSON values as the path argument to probe schema validation",
         9, False),
        ("S8", "Exploit Variant Enumeration",
         "Mutate the confirmed bypass to assess patch completeness and coverage",
         8, False),
    ]
    y = 2.0
    for num, name, desc, n, hi in inference:
        _scenario_cube(s, 7.1, y, num, name, desc, n, highlight=hi)
        y += 1.08

    add_textbox(s, 0.83, 6.3, 11.6, 0.38,
                "Each test scenario's observed response directly informed the design of the next. "
                "S6 (red) is the scenario in which a path access succeeded contrary to expectation.",
                size_pt=12, color=DARK_GRAY)


def _tool_col_header(slide, x, y, tool_name, n_inputs):
    """Column header for a tool split."""
    add_rect(slide, x, y, 5.55, 0.04, fill_color=RED)
    add_textbox(slide, x, y + 0.08, 3.5, 0.36,
                tool_name, size_pt=15, bold=True, color=WHITE)
    add_textbox(slide, x + 3.6, y + 0.08, 1.9, 0.36,
                f"{n_inputs} inputs", size_pt=12, color=DARK_GRAY, align="right")


def _tool_row(slide, x, y, input_ex, output_ex):
    """One input/output example row inside a tool column."""
    add_textbox(slide, x, y, 5.45, 0.32,
                input_ex, size_pt=12, color=LIGHT_GRAY, italic=True)
    add_textbox(slide, x + 0.18, y + 0.32, 5.27, 0.3,
                output_ex, size_pt=11, color=DARK_GRAY)


def slide_06_phase1(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "S1 — Baseline Characterisation: Nominal Server Behaviour")

    add_textbox(s, 0.83, 1.2, 1.7, 0.42,
                "13 test cases", size_pt=14, bold=True, color=RED)
    add_textbox(s, 2.6, 1.2, 10.0, 0.42,
                "All inputs are expected to succeed. "
                "Objective: establish the nominal response envelope so that "
                "any subsequent deviation constitutes an observable anomaly.",
                size_pt=14, color=WHITE)

    divider(s, 1.75)

    # ── read_file column  (10 test cases) ───────────────────────────────────
    _tool_col_header(s, 0.83, 1.83, "read_file", 10)

    rf_rows = [
        ("/sandbox/readme.txt",
         "Response: full plaintext content — nominal baseline"),
        ("/sandbox/config.txt",
         "Response: key=value pairs verbatim — no server-side transformation"),
        ("/sandbox/secret.txt",
         "Response: content returned — filename semantics not inspected by server"),
        ("/sandbox/empty.txt  [0 bytes]",
         "Response: empty string body, isError: false — zero-byte files valid"),
        ("/sandbox/.hidden",
         "Response: content returned — POSIX hidden files treated identically"),
        ("/sandbox/large.txt  [500 lines ~50 KB]",
         "Response: full content, no truncation — no implicit size limit observed"),
        ("/sandbox/subdir/nested.txt",
         "Response: content returned — single-level subdirectory traversal works"),
        ("/sandbox/subdir/deep/deep.txt",
         "Response: content returned — two-level nesting also valid"),
        ("/sandbox/symlink.txt  [symlink → readme.txt]",
         "Response: target content — server follows symlinks transparently"),
        ("/sandbox/file with spaces.txt",
         "Response: content returned — whitespace in filenames not mishandled"),
    ]
    y = 2.37
    for inp, out in rf_rows:
        _tool_row(s, 0.83, y, inp, out)
        y += 0.44

    # ── list_directory column  (3 test cases) ───────────────────────────────
    _tool_col_header(s, 7.1, 1.83, "list_directory", 3)

    ld_rows = [
        ("/sandbox/",
         "Response: [FILE] readme.txt  [FILE] config.txt  [DIR] subdir  …\n"
         "All entries present, consistent [FILE]/[DIR] prefix notation"),
        ("/sandbox/subdir/",
         "Response: [DIR] deep  [FILE] nested.txt\n"
         "Subdirectory listing structurally identical to root listing"),
        ("/sandbox/empty_dir/",
         "Response: empty body, isError: false\n"
         "Empty directory returns no entries, not an error condition"),
    ]
    y = 2.37
    for inp, out in ld_rows:
        _tool_row(s, 7.1, y, inp, out)
        y += 1.36

    divider(s, 6.72, color=DARK_GRAY)
    add_textbox(s, 0.83, 6.8, 11.6, 0.42,
                "Finding: server behaviour is deterministic and consistent across all 13 cases. "
                "Critically, error messages expose the fully resolved absolute path — "
                "a side-channel exploited in subsequent scenarios.",
                size_pt=13, bold=True, color=WHITE)


def slide_07_phases234(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "S2–S4 — Error Analysis, Path Resolution & Boundary Enforcement")

    # table column headers
    divider(s, 1.2)
    add_textbox(s, 0.83, 1.25, 2.15, 0.32, "Test Scenario",
                size_pt=12, bold=True, color=DARK_GRAY)
    add_textbox(s, 3.05, 1.25, 4.1, 0.32, "read_file  — input / response",
                size_pt=12, bold=True, color=RED)
    add_textbox(s, 7.25, 1.25, 4.1, 0.32, "list_directory  — input / response",
                size_pt=12, bold=True, color=RED)
    add_textbox(s, 11.45, 1.25, 0.9, 0.32, "n",
                size_pt=12, bold=True, color=DARK_GRAY, align="right")
    divider(s, 1.62, color=DARK_GRAY)

    # ── S2 — Error Response Analysis ────────────────────────────────────────
    y = 1.68
    add_textbox(s, 0.83, y, 2.15, 0.32,
                "S2  Error Response\nAnalysis",
                size_pt=12, bold=True, color=WHITE)
    add_textbox(s, 0.83, y + 0.65, 2.1, 0.55,
                "Map error structure for missing, denied, and out-of-scope paths",
                size_pt=10, color=LIGHT_GRAY)

    add_textbox(s, 3.05, y, 4.1, 0.28,
                "/sandbox/nonexistent.txt", size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 3.05, y + 0.28, 4.1, 0.25,
                "ENOENT: no such file — full resolved path included in body",
                size_pt=10, color=DARK_GRAY)
    add_textbox(s, 3.05, y + 0.56, 4.1, 0.25,
                "/sandbox/subdir/missing.txt", size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 3.05, y + 0.81, 4.1, 0.25,
                "ENOENT — nested missing path, same error envelope",
                size_pt=10, color=DARK_GRAY)

    add_textbox(s, 7.25, y, 4.1, 0.28,
                "/sandbox/absent_dir/", size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 7.25, y + 0.28, 4.1, 0.25,
                "ENOENT — identical error structure to read_file",
                size_pt=10, color=DARK_GRAY)
    add_textbox(s, 7.25, y + 0.56, 4.1, 0.28,
                "/outside/path/",
                size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 7.25, y + 0.84, 4.1, 0.25,
                "Access denied — resolved path leaked in error body",
                size_pt=10, color=DARK_GRAY)

    add_textbox(s, 11.45, y, 0.9, 0.32, "8",
                size_pt=13, bold=True, color=DARK_GRAY, align="right")
    add_textbox(s, 0.83, y + 1.12, 11.6, 0.28,
                "Observation: error messages consistently expose the canonical resolved path — "
                "a reliable side-channel for inferring server-side path resolution logic.",
                size_pt=11, bold=True, color=WHITE)

    divider(s, y + 1.48, color=MID_GRAY)

    # ── S3 — Path Resolution Behaviour ──────────────────────────────────────
    y = y + 1.56
    add_textbox(s, 0.83, y, 2.15, 0.32,
                "S3  Path Resolution\nBehaviour",
                size_pt=12, bold=True, color=WHITE)
    add_textbox(s, 0.83, y + 0.65, 2.1, 0.55,
                "Determine whether ACL evaluation operates on canonical or raw form",
                size_pt=10, color=LIGHT_GRAY)

    add_textbox(s, 3.05, y, 4.1, 0.28,
                "/sandbox/./././readme.txt", size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 3.05, y + 0.28, 4.1, 0.25,
                "Content returned — redundant current-dir refs collapsed",
                size_pt=10, color=DARK_GRAY)
    add_textbox(s, 3.05, y + 0.56, 4.1, 0.25,
                "//sandbox//readme.txt", size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 3.05, y + 0.81, 4.1, 0.25,
                "Content returned — duplicate separators normalised",
                size_pt=10, color=DARK_GRAY)

    add_textbox(s, 7.25, y, 4.1, 0.28,
                "/sandbox/subdir/deep/../", size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 7.25, y + 0.28, 4.1, 0.25,
                "Lists subdir/ — parent-dir component resolved correctly",
                size_pt=10, color=DARK_GRAY)
    add_textbox(s, 7.25, y + 0.56, 4.1, 0.25,
                "///sandbox///", size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 7.25, y + 0.81, 4.1, 0.25,
                "Root listing returned — triple slashes collapsed",
                size_pt=10, color=DARK_GRAY)

    add_textbox(s, 11.45, y, 0.9, 0.32, "8",
                size_pt=13, bold=True, color=DARK_GRAY, align="right")
    add_textbox(s, 0.83, y + 1.12, 11.6, 0.28,
                "Observation: normalisation precedes I/O. "
                "Key question raised: does ACL evaluation use the normalised path or the raw input string?",
                size_pt=11, bold=True, color=WHITE)

    divider(s, y + 1.48, color=MID_GRAY)

    # ── S4 — Malformed Input Handling ───────────────────────────────────────
    y = y + 1.56
    add_textbox(s, 0.83, y, 2.15, 0.32,
                "S4  Malformed Input\nHandling",
                size_pt=12, bold=True, color=WHITE)
    add_textbox(s, 0.83, y + 0.65, 2.1, 0.55,
                "Assess server response to syntactically invalid path strings",
                size_pt=10, color=LIGHT_GRAY)

    add_textbox(s, 3.05, y, 4.1, 0.65,
                '"aaaaaaa"  "12345"  "???"  "\\\\"  "   "  "\\n"  "\\t"',
                size_pt=11, color=LIGHT_GRAY, italic=True)
    add_textbox(s, 3.05, y + 0.65, 4.1, 0.25,
                "All 14 inputs: Access denied — path outside allowed directories",
                size_pt=10, color=DARK_GRAY)
    add_textbox(s, 3.05, y + 0.9, 4.1, 0.25,
                "Resolved form exposed in each error (e.g. /mnt/\"aaaaaaa\")",
                size_pt=10, color=DARK_GRAY)

    add_textbox(s, 7.25, y, 4.1, 0.32,
                "(read_file only — S4 targets file-read path)", size_pt=11, color=MID_GRAY)
    add_textbox(s, 7.25, y + 0.35, 4.1, 0.28,
                "list_directory excluded: malformed string inputs are\n"
                "a file-access concern, not a directory enumeration concern",
                size_pt=10, color=MID_GRAY)

    add_textbox(s, 11.45, y, 0.9, 0.32, "14",
                size_pt=13, bold=True, color=DARK_GRAY, align="right")
    add_textbox(s, 0.83, y + 1.2, 11.6, 0.28,
                "Observation: access-control boundary is consistently enforced. "
                "Error output confirms the server resolves the input before rejecting — "
                "boundary check operates on the canonical path, not the raw string.",
                size_pt=11, bold=True, color=WHITE)


def slide_08_phase5(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "S5 — Near-Boundary Input Testing: First Inter-Version Divergence")

    add_textbox(s, 0.83, 1.2, 2.0, 0.4,
                "14 test cases", size_pt=14, bold=True, color=RED)
    add_textbox(s, 2.9, 1.2, 9.7, 0.4,
                "Inputs that are syntactically plausible but semantically ambiguous — "
                "targeting edge cases at the boundary of the access-control policy.",
                size_pt=14, color=WHITE)

    divider(s, 1.7)

    # column headers
    add_textbox(s, 0.83, 1.76, 2.1, 0.3,
                "Input class", size_pt=12, bold=True, color=DARK_GRAY)
    add_textbox(s, 3.05, 1.76, 4.2, 0.3,
                "read_file  (12 test cases)", size_pt=12, bold=True, color=RED)
    add_textbox(s, 7.35, 1.76, 4.7, 0.3,
                "list_directory  (2 test cases)", size_pt=12, bold=True, color=RED)
    divider(s, 2.12, color=DARK_GRAY)

    # ── read_file test cases ─────────────────────────────────────────────────
    rf = [
        ("Windows path separator",
         r"/sandbox\readme.txt",
         "v2025.3.28: content returned",
         "v2025.7.29: Access denied  [patch effective]",
         True),
        ("Null-byte injection",
         "/sandbox/readme.txt\\x00.png",
         "v2025.3.28: content returned",
         "v2025.7.29: Access denied  [patch effective]",
         True),
        ("Parent-directory traversal",
         "../../sandbox/readme.txt",
         "Both versions: Access denied",
         "Canonical form resolves outside the allowed root",
         False),
        ("Unicode fraction slash U+2215",
         "/sandbox\u2215readme.txt",
         "Both versions: Access denied",
         "Non-ASCII separator not decoded — path unrecognised",
         False),
        ("URL-encoded separator %2F",
         "/sandbox%2Freadme.txt",
         "Both versions: Access denied",
         "Percent-encoding not decoded — treated as literal filename",
         False),
        ("Pathological redundancy (300 chars)",
         "/sandbox/a/../a/../a/… (truncated)",
         "Both versions: content returned",
         "Responses differ in body text — unexpected divergence",
         False),
    ]

    y = 2.2
    for label, inp, r1, r2, is_patch in rf:
        lc = RED if is_patch else WHITE
        add_textbox(s, 0.83, y, 2.1, 0.27, label, size_pt=11, bold=True, color=lc)
        add_textbox(s, 3.05, y, 4.2, 0.25, inp,
                    size_pt=10, color=LIGHT_GRAY, italic=True)
        rc = RED if is_patch else DARK_GRAY
        add_textbox(s, 3.05, y + 0.25, 4.2, 0.22, r1, size_pt=10, color=DARK_GRAY)
        add_textbox(s, 3.05, y + 0.47, 4.2, 0.22, r2, size_pt=10, color=rc, bold=is_patch)
        y += 0.75

    # ── list_directory test cases ────────────────────────────────────────────
    ld = [
        ("Windows separator on directory path",
         r"/sandbox\subdir\ ",
         "v2025.3.28: directory listed",
         "v2025.7.29: Access denied  [patch effective]",
         True),
        ("Duplicate leading separators",
         "//sandbox//",
         "Both versions: root listing returned",
         "Duplicate separators normalised — consistent with S3",
         False),
    ]
    y_ld = 2.2
    for label, inp, r1, r2, is_patch in ld:
        lc = RED if is_patch else WHITE
        add_textbox(s, 7.35, y_ld, 4.7, 0.27, label, size_pt=11, bold=True, color=lc)
        add_textbox(s, 7.35, y_ld + 0.27, 4.7, 0.25, inp,
                    size_pt=10, color=LIGHT_GRAY, italic=True)
        rc = RED if is_patch else DARK_GRAY
        add_textbox(s, 7.35, y_ld + 0.52, 4.7, 0.22, r1, size_pt=10, color=DARK_GRAY)
        add_textbox(s, 7.35, y_ld + 0.74, 4.7, 0.22, r2, size_pt=10, color=rc, bold=is_patch)
        y_ld += 1.35

    divider(s, 6.73, color=DARK_GRAY)
    add_textbox(s, 0.83, 6.82, 7.5, 0.38,
                "3 of 14 exhibit inter-version divergence. 11 of 14 remain identical. "
                "The access-control boundary appears structurally sound.",
                size_pt=12, color=WHITE)
    add_textbox(s, 8.4, 6.82, 4.1, 0.38,
                "Hypothesis revised: the boundary logic itself may be flawed.",
                size_pt=12, bold=True, color=RED)


def slide_09_phase6_discovery(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "Phase 6 — The Discovery: Prefix Bypass")

    add_textbox(s, 0.83, 1.25, 11.6, 0.45,
                "I stopped using traversal dots. Instead I asked: "
                "what if a path simply STARTS with the sandbox name?",
                size_pt=16, color=WHITE)

    # The hypothesis vs. actual
    add_textbox(s, 0.83, 1.85, 5.5, 0.38, "My hypothesis",
                size_pt=14, bold=True, color=RED)
    add_textbox(s, 0.83, 2.25, 5.5, 0.85,
                "The server checks: does path start with /tmp/mcp_sandbox?\n"
                "If so — allowed. /tmp/mcp_sandbox_escape also starts with\n"
                "that prefix. Would it accept it?",
                size_pt=14, color=LIGHT_GRAY, italic=True)

    add_textbox(s, 6.83, 1.85, 5.5, 0.38, "What happened",
                size_pt=14, bold=True, color=RED)
    code_block(s, 6.83, 2.25, 5.5, 1.35, [
        ("path: /tmp/mcp_sandbox_escape/secret.txt", CODE_GREEN),
        ("",                                          WHITE),
        ("OLD → returned content",                    CODE_GREEN),
        ("NEW → Access denied",                       RGBColor(0xFF, 0x60, 0x60)),
    ])

    add_textbox(s, 0.83, 3.7, 11.6, 0.38,
                "It worked on the old version. Something succeeded when it should have failed.",
                size_pt=16, bold=True, color=RED)

    divider(s, 4.2, color=DARK_GRAY)

    variants = [
        ("/tmp/mcp_sandbox_escape/file1.txt",   "leaked"),
        ("/tmp/mcp_sandbox_escape/file2.txt",   "leaked"),
        ("/tmp/mcp_sandboxevil/secret.txt",     "leaked  (no separator needed)"),
        ("list /tmp/mcp_sandbox_escape/",       "directory listing leaked"),
    ]
    y = 4.35
    add_textbox(s, 0.83, y, 11.6, 0.35,
                "Confirmed with 6 variants — all leaked on old version:",
                size_pt=14, color=WHITE)
    y += 0.4
    for path, result in variants:
        add_textbox(s, 1.1,  y, 7.5, 0.35, path,   size_pt=13, color=LIGHT_GRAY, italic=True)
        add_textbox(s, 8.7,  y, 3.6, 0.35, result, size_pt=13, bold=True, color=RED)
        y += 0.38


def slide_10_cve(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "The CVE: What Prefix Bypass Means")

    add_textbox(s, 0.83, 1.25, 11.6, 0.45,
                "The old server validated the path with a naive string-prefix check:",
                size_pt=16, color=LIGHT_GRAY)

    code_block(s, 0.83, 1.8, 11.6, 1.05, [
        ("// Vulnerable check (conceptual)",                              DARK_GRAY),
        ("if (!resolvedPath.startsWith(allowedRoot)) throw AccessDenied", RGBColor(0xFF, 0x80, 0x80)),
        ("",                                                               WHITE),
        ("// /tmp/mcp_sandbox_escape  starts with  /tmp/mcp_sandbox  → passes!", CODE_YELLOW),
    ])

    code_block(s, 0.83, 3.0, 11.6, 0.95, [
        ("// Patched check (v2025.7.29)",                                  DARK_GRAY),
        ("if (!resolvedPath.startsWith(allowedRoot + path.sep)) throw …", CODE_GREEN),
        ("",                                                               WHITE),
        ("// Forces a separator after the sandbox name → sibling dirs blocked", CODE_YELLOW),
    ])

    rows = [
        ("Impact",   "agent can read any file in any directory whose name "
                     "shares a prefix with the sandbox"),
        ("Scope",    "any MCP filesystem server using startsWith() without a trailing separator"),
        ("Fix",      "append path.sep (or '/') to the allowed root before the check"),
    ]
    y = 4.15
    for label, detail in rows:
        add_textbox(s, 0.83, y, 1.5, 0.42, label, size_pt=15, bold=True, color=RED)
        add_textbox(s, 2.4,  y, 10.1, 0.5, detail, size_pt=14, color=WHITE)
        y += 0.55


def slide_11_phase7(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "Phase 7 — Type Confusion: A Side Finding")

    add_textbox(s, 0.83, 1.25, 11.6, 0.45,
                "After the prefix bypass I wondered: what if 'path' is not a string at all? "
                "9 probes with wrong JSON types.",
                size_pt=16, color=WHITE)

    examples = [
        ("path: 42",    "unhandled exception",        "structured error"),
        ("path: null",  "ENOENT on literal 'null'",   "schema validation error"),
        ("path: true",  "reads 'true' as path string","blocked before I/O"),
        ("path: []",    "TypeError in handler",       "structured error"),
        ("path: {}",    "TypeError in handler",       "structured error"),
    ]

    add_textbox(s, 1.1,  2.1, 2.8, 0.38, "Input",     size_pt=14, bold=True, color=RED)
    add_textbox(s, 4.0,  2.1, 3.9, 0.38, "Old behaviour", size_pt=14, bold=True, color=RED)
    add_textbox(s, 7.95, 2.1, 4.5, 0.38, "New behaviour", size_pt=14, bold=True, color=RED)
    divider(s, 2.55)

    y = 2.65
    for inp, old, new in examples:
        add_textbox(s, 1.1,  y, 2.8, 0.4, inp, size_pt=13, color=LIGHT_GRAY, italic=True)
        add_textbox(s, 4.0,  y, 3.9, 0.4, old, size_pt=13, color=LIGHT_GRAY)
        add_textbox(s, 7.95, y, 4.5, 0.4, new, size_pt=13, color=WHITE)
        y += 0.45

    divider(s, y + 0.1, color=DARK_GRAY)
    add_textbox(s, 0.83, y + 0.25, 11.6, 0.55,
                "Not exploitable — but a caller can fingerprint the server version "
                "just by sending path: 42 and reading the error shape.",
                size_pt=15, color=LIGHT_GRAY, italic=True)


def slide_12_phase8(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "Phase 8 — Follow-Up: Confirming the CVE Family")

    add_textbox(s, 0.83, 1.25, 11.6, 0.45,
                "I now know the prefix bypass works. Phase 8 tests whether the "
                "patch is a real fix or just one blacklisted spelling.",
                size_pt=16, color=WHITE)

    mutations = [
        ("bypass + trailing slash",    "/tmp/mcp_sandbox_escape/",       "OLD leaked  →  NEW blocked"),
        ("bypass + double slash",      "/tmp/mcp_sandbox_escape//",      "OLD leaked  →  NEW blocked"),
        ("different suffix _extra",    "/tmp/mcp_sandbox_extra/file",    "OLD leaked  →  NEW blocked"),
        ("dot in escape path",         "/tmp/mcp_sandbox_escape/./file", "OLD leaked  →  NEW blocked"),
        ("list escape dir",            "list_directory on escape dir",   "OLD listed  →  NEW blocked"),
        ("list via double slash",      "list_directory + //",            "OLD listed  →  NEW blocked"),
    ]

    y = 2.0
    for label, path, result in mutations:
        add_textbox(s, 0.83, y, 3.5, 0.4, label,  size_pt=14, bold=True, color=WHITE)
        add_textbox(s, 4.4,  y, 4.6, 0.4, path,   size_pt=13, color=LIGHT_GRAY, italic=True)
        add_textbox(s, 9.1,  y, 3.4, 0.4, result, size_pt=13, bold=True, color=RED)
        y += 0.47

    divider(s, y + 0.1, color=DARK_GRAY)
    add_textbox(s, 0.83, y + 0.25, 11.6, 0.55,
                "Every variant blocked by the patched version. "
                "The fix is real path resolution — not a one-liner blacklist.",
                size_pt=15, bold=True, color=WHITE)


def slide_13_bonus_versions(prs):
    s = add_slide(prs, LO_TITLE_ONLY_BLACK)
    set_title(s, "Bonus: Version Comparison Validates the Find")

    add_textbox(s, 0.83, 1.25, 11.6, 0.45,
                "Running the same 90 probes against v2025.3.28 and v2025.7.29 in parallel "
                "confirmed the CVE independently.",
                size_pt=16, color=LIGHT_GRAY)

    stats = [
        ("90",  "total probes"),
        ("43",  "identical\n(phases 1–4)"),
        ("23",  "diverged\n(phases 5–8)"),
        ("15",  "PATCH\nEFFECTIVE"),
        ("8",   "type confusion\ndivergence"),
    ]
    x = 0.83
    for big, small in stats:
        add_textbox(s, x, 2.1, 2.2, 1.0, big,   size_pt=46, bold=True, color=RED, align="center")
        add_textbox(s, x, 3.1, 2.2, 0.7, small, size_pt=12, color=LIGHT_GRAY, align="center")
        x += 2.35

    divider(s, 4.05, color=DARK_GRAY)

    add_textbox(s, 0.83, 4.2, 11.6, 0.45,
                "Phases 1–4 identical → patch did not regress benign behaviour.",
                size_pt=15, color=WHITE)
    add_textbox(s, 0.83, 4.7, 11.6, 0.45,
                "15 PATCH EFFECTIVE rows = every prefix-bypass variant "
                "blocked by v2025.7.29 and leaked by v2025.3.28.",
                size_pt=15, color=WHITE)
    add_textbox(s, 0.83, 5.2, 11.6, 0.45,
                "The probing methodology was sensitive enough to catch a real CVE "
                "and confirm the fix — without reading the source code.",
                size_pt=15, bold=True, color=RED)


def slide_14_closing(prs):
    s = add_slide(prs, LO_CLOSING_BLACK)
    add_textbox(s, 1.5, 2.3, 10.33, 0.9,
                "Key Takeaway",
                size_pt=32, bold=True, color=WHITE, align="center")
    add_textbox(s, 1.5, 3.3, 10.33, 1.1,
                "Systematic hypothesis-first probing of a black-box MCP server\n"
                "— starting from baseline, reacting to every output —\n"
                "was enough to discover a path-validation CVE independently.",
                size_pt=18, color=LIGHT_GRAY, align="center")
    add_textbox(s, 1.5, 4.55, 10.33, 0.5,
                "This methodology becomes the empirical foundation for the "
                "thesis risk-scoring framework.",
                size_pt=15, color=RED, align="center")
    add_textbox(s, 1.5, 5.2, 10.33, 0.45,
                "Questions?",
                size_pt=22, bold=True, color=WHITE, align="center")


# ---------------------------------------------------------------------------
# Slide removal helper
# ---------------------------------------------------------------------------

def remove_original_slides(prs: Presentation, keep_from: int) -> None:
    prs_part = prs.part
    prs_el = prs_part._element
    sld_id_lst = prs_el.find(qn("p:sldIdLst"))
    all_sld_ids = list(sld_id_lst)
    for sld_id in all_sld_ids[:keep_from]:
        r_id = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs_part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

BUILDERS = [
    slide_01_title,
    slide_02_what_is_mcp,
    slide_03_json_protocol,
    slide_04_two_tools,
    slide_05_methodology,
    slide_06_phase1,
    slide_07_phases234,
    slide_08_phase5,
    slide_09_phase6_discovery,
    slide_10_cve,
    slide_11_phase7,
    slide_12_phase8,
    slide_13_bonus_versions,
    slide_14_closing,
]


def build():
    print(f"Loading template: {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))
    original_count = len(prs.slides)
    print(f"Template has {original_count} existing slides — will replace them.")

    for fn in BUILDERS:
        fn(prs)
        print(f"  + {fn.__name__}")

    remove_original_slides(prs, original_count)
    print(f"Removed {original_count} original slides.")

    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
