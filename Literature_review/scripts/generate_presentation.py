"""Generate thesis seminar presentation for Tomer Ovadya.

Creates a PowerPoint with:
  Part 1 (3 min): Title, Problem & Motivation, Research Goals
  Part 2 (5 min): Domains Overview, Threats, Defenses, Scoring, Comparison Table
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

# ── Constants ────────────────────────────────────────────────────────────────
TEAL = RGBColor(0x00, 0x9D, 0xBB)
DARK_TEAL = RGBColor(0x00, 0x61, 0x89)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x44, 0x54, 0x6A)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE = RGBColor(0xF4, 0x90, 0x21)
CHECK_GREEN = RGBColor(0x00, 0x7A, 0x33)
RED_GAP = RGBColor(0xCC, 0x00, 0x00)
LIGHT_GREEN = RGBColor(0xE0, 0xF7, 0xE0)

HEADING = "Calibri Light"
BODY = "Calibri"

TEMPLATE = Path(r"C:/Users/user/Downloads/Tomerovadya (2).pptx")
OUTPUT = Path(r"C:/Users/user/Documents/GitHub/MCP/Literature_review/TomerOvadya.pptx")


# ── Helpers ──────────────────────────────────────────────────────────────────
def delete_all_slides(prs):
    """Remove every existing slide from the presentation."""
    for sld_id in list(prs.slides._sldIdLst):
        r_id = sld_id.get(qn("r:id"))
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(sld_id)


def tb(slide, left, top, width, height):
    """Add a word-wrapped textbox; return its text_frame."""
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    return tf


def run(para, text, font=BODY, size=18, color=BLACK, bold=False):
    """Append a styled run to *para*."""
    r = para.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return r


def slide_title(slide, text):
    """Standard teal title at the top of a content slide."""
    tf = tb(slide, Inches(1.0), Inches(0.3), Inches(11.3), Inches(0.8))
    run(tf.paragraphs[0], text, HEADING, 36, TEAL, bold=True)


def new_slide(prs, layout):
    """Add slide, silence placeholder text."""
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        try:
            if ph.has_text_frame:
                for p in ph.text_frame.paragraphs:
                    p.clear()
        except Exception:
            pass
    return slide


# ── Build ────────────────────────────────────────────────────────────────────
prs = Presentation(str(TEMPLATE))
delete_all_slides(prs)
layout = prs.slide_layouts[11]  # branded custom layout with background image


# ═══════════════════════════════════════
# PART 1 — Research Topic Introduction
# ═══════════════════════════════════════

# ── Slide 1: Title ───────────────────
s = new_slide(prs, layout)

tf = tb(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.0))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "Dynamic Risk Scoring for\nMCP Agent Access", HEADING, 40, TEAL, bold=True)

tf2 = tb(s, Inches(1.5), Inches(4.0), Inches(10.3), Inches(2.0))
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "Tomer Ovadya", BODY, 24, BLACK, bold=True)

p = tf2.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(12)
run(p, "Advisor: Prof. Asaf Shabtai", BODY, 20, DARK_GRAY)

p = tf2.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)
run(p, "Ben-Gurion University of the Negev", BODY, 18, DARK_GRAY)


# ── Slide 2: Problem & Motivation ────
s = new_slide(prs, layout)
slide_title(s, "Problem & Motivation")

tf = tb(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(5.5))

bullets = [
    ("MCP (Model Context Protocol) enables AI agents to access external tools and services", False),
    ("Rapid adoption: 67,000+ public servers; integrated by Anthropic, OpenAI, Google", False),
    ("Current security approach: Binary allow / deny decisions", False),
    ("No mechanism exists to quantify the risk level of agent requests", True),
    ("Need: Dynamic, context-aware risk assessment for trusted agent access", True),
]
for i, (text, bold) in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(14)
    run(p, f"\u2022  {text}", BODY, 20, BLACK, bold)


# ── Slide 3: Research Goals ──────────
s = new_slide(prs, layout)
slide_title(s, "Research Goals")

tf = tb(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(5.5))

p = tf.paragraphs[0]
run(p, "Research Question", HEADING, 24, TEAL, bold=True)

p = tf.add_paragraph()
p.space_before = Pt(8)
run(p, "How can we determine if an MCP agent is trusted,", BODY, 22, DARK_GRAY, bold=True)

p = tf.add_paragraph()
run(p, "and calculate a dynamic risk score for its access request?", BODY, 22, DARK_GRAY, bold=True)

p = tf.add_paragraph()
p.space_before = Pt(30)
run(p, "Goals", HEADING, 24, TEAL, bold=True)

for g in [
    "Design a dynamic risk scoring system (1\u201310 scale)",
    "Evaluate agent access requests using multiple risk factors",
    "Provide real-time, context-aware risk assessment",
    "Enable granular access decisions \u2014 beyond binary allow / deny",
]:
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    run(p, f"\u2022  {g}", BODY, 20, BLACK)


# ═══════════════════════════════════════
# PART 2 — Literature Review
# ═══════════════════════════════════════

# ── Slide 4: Literature Domains ──────
s = new_slide(prs, layout)
slide_title(s, "Literature Review \u2014 Key Domains")

domains = [
    ("1. MCP Security &\n    Threat Landscape",
     "Threat taxonomies \u00b7 Protocol analysis\nEcosystem vulnerabilities",
     TEAL),
    ("2. Defense Frameworks\n    & Access Control",
     "Sandbox probing \u00b7 Cascaded defense\nPrivilege control \u00b7 OS permissions",
     DARK_TEAL),
    ("3. Attack Methods &\n    Tool Poisoning",
     "Tool poisoning benchmarks\nParasitic attacks \u00b7 Implicit poisoning",
     DARK_TEAL),
    ("4. Risk & Trust\n    Scoring",
     "Vulnerability scoring \u00b7 TRiSM framework\nGuard agents \u00b7 Step-level guardrails",
     TEAL),
]

box_w, box_h = Inches(5.2), Inches(2.3)
gap_x, gap_y = Inches(0.6), Inches(0.4)
x0, y0 = Inches(1.0), Inches(1.6)

for i, (title, desc, bg) in enumerate(domains):
    row, col = divmod(i, 2)
    x = x0 + col * (box_w + gap_x)
    y = y0 + row * (box_h + gap_y)

    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()

    stf = shape.text_frame
    stf.word_wrap = True
    stf.margin_left = Inches(0.2)
    stf.margin_right = Inches(0.2)
    stf.margin_top = Inches(0.25)

    p = stf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, title, HEADING, 22, WHITE, bold=True)

    p = stf.add_paragraph()
    p.space_before = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    run(p, desc, BODY, 16, WHITE)


# ── Slide 5: MCP Threats & Attack Methods ──
s = new_slide(prs, layout)
slide_title(s, "MCP Threats & Attack Methods")

tf = tb(s, Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.5))

papers_threats = [
    ("MCP Landscape", "Hou et al., 2025",
     "First comprehensive threat taxonomy \u2014 4 attacker types, 16 scenarios"),
    ("When MCP Servers Attack", "Zhao et al., 2025",
     "Malicious server analysis \u2014 12 attack categories across 6 components"),
    ("MCPTox", "Wang et al., 2025",
     "First tool poisoning benchmark on real MCP servers"),
    ("MCP-ITP", "Li et al., 2026",
     "Automated implicit tool poisoning \u2014 <1% detection rate"),
    ("Breaking the Protocol", "Maloyan & Namiot, 2026",
     "Architectural flaws in MCP spec itself, not implementation bugs"),
]

for i, (name, authors, desc) in enumerate(papers_threats):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(10) if i > 0 else Pt(0)
    run(p, f"{name} ", BODY, 18, TEAL, bold=True)
    run(p, f"({authors})", BODY, 14, DARK_GRAY)

    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run(p, f"    \u2192 {desc}", BODY, 16, BLACK)

p = tf.add_paragraph()
p.space_before = Pt(20)
run(p, "Gap: ", BODY, 18, ORANGE, bold=True)
run(p, "Focus on identifying threats \u2014 not on quantifying risk", BODY, 18, DARK_GRAY, bold=True)


# ── Slide 6: Defense Frameworks ──────
s = new_slide(prs, layout)
slide_title(s, "Defense Frameworks & Access Control")

tf = tb(s, Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.5))

papers_defense = [
    ("MCPShield", "Zhou et al., 2026",
     "Trust calibration via sandbox probing",
     "Self-hosted only, no risk score"),
    ("MCP-Guard", "Xing et al., 2026",
     "Cascaded defense: regex \u2192 neural \u2192 LLM",
     "Binary classification only"),
    ("Progent", "Shi et al., 2025",
     "Programmable privilege control for agents",
     "Manual policies, no scoring"),
    ("AgentBound", "Buhler et al., 2025",
     "OS-level permissions for MCP servers",
     "Binary allow/deny, no context"),
]

for i, (name, authors, desc, gap) in enumerate(papers_defense):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(12) if i > 0 else Pt(0)
    run(p, f"{name} ", BODY, 18, TEAL, bold=True)
    run(p, f"({authors})", BODY, 14, DARK_GRAY)

    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run(p, f"    \u2192 {desc}", BODY, 16, BLACK)

    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run(p, f"    \u2717 {gap}", BODY, 15, RED_GAP)

p = tf.add_paragraph()
p.space_before = Pt(16)
run(p, "Gap: ", BODY, 18, ORANGE, bold=True)
run(p, "Binary decisions \u2014 no quantitative risk assessment", BODY, 18, DARK_GRAY, bold=True)


# ── Slide 7: Risk & Trust Scoring ────
s = new_slide(prs, layout)
slide_title(s, "Risk & Trust Scoring Approaches")

tf = tb(s, Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.5))

papers_risk = [
    ("Desc. to Score", "Jafarikhah et al., 2026",
     "LLMs predict CVSS vulnerability scores",
     "Not MCP-specific, no agent context"),
    ("TRiSM Survey", "Raza et al., 2025",
     "Trust, Risk & Security framework for multi-agent systems",
     "Survey only \u2014 no implementation"),
    ("GuardAgent", "Xiang et al., 2025",
     "Knowledge-guided guard agent with reasoning",
     "Binary allow/deny, no risk scoring"),
    ("ToolSafe", "Mou et al., 2026",
     "Step-level safety guardrails with proactive feedback",
     "3-level scale only, no continuous monitoring"),
]

for i, (name, authors, desc, gap) in enumerate(papers_risk):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(12) if i > 0 else Pt(0)
    run(p, f"{name} ", BODY, 18, TEAL, bold=True)
    run(p, f"({authors})", BODY, 14, DARK_GRAY)

    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run(p, f"    \u2192 {desc}", BODY, 16, BLACK)

    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run(p, f"    \u2717 {gap}", BODY, 15, RED_GAP)

p = tf.add_paragraph()
p.space_before = Pt(16)
run(p, "Gap: ", BODY, 18, ORANGE, bold=True)
run(p, "No dynamic, MCP-specific risk scoring system exists", BODY, 18, DARK_GRAY, bold=True)


# ── Slide 8: Comparison Table ────────
s = new_slide(prs, layout)
slide_title(s, "Research Gap Analysis")

HEADERS = [
    "Paper", "MCP\nSpecific", "Risk\nScore", "Agent\nTrust",
    "Real-\nTime", "Tool-\nLevel", "Context-\nAware", "Continuous\nMonitoring",
]

# (name, mcp, score, trust, rt, tool, ctx, cont)
TABLE_DATA = [
    ("MCPShield (Zhou et al.)",      "\u2713", "\u2014", "\u2713", "\u2014", "\u2713", "\u2713", "\u2713"),
    ("MCP-Guard (Xing et al.)",      "\u2713", "\u2014", "\u2014", "\u2713", "\u2713", "\u2014", "\u2014"),
    ("Progent (Shi et al.)",         "\u2014", "\u2014", "\u2014", "\u2713", "\u2713", "\u2713", "\u2014"),
    ("AgentBound (Buhler et al.)",   "\u2713", "\u2014", "\u2014", "\u2713", "\u2713", "\u2014", "\u2014"),
    ("Desc. to Score (Jafarikhah)",  "\u2014", "\u2713", "\u2014", "\u2014", "\u2014", "\u2014", "\u2014"),
    ("MindGuard (Wang et al.)",      "\u2014", "\u2014", "\u2014", "\u2713", "\u2713", "\u2713", "\u2014"),
    ("TRiSM (Raza et al.)",         "\u2014", "\u2014", "\u2713", "\u2014", "\u2014", "\u2713", "\u2014"),
    ("GuardAgent (Xiang et al.)",    "\u2014", "\u2014", "\u2014", "\u2713", "\u2713", "\u2713", "\u2014"),
    ("ToolSafe (Mou et al.)",        "\u2014", "\u2014", "\u2014", "\u2713", "\u2713", "\u2014", "\u2014"),
    ("Our Approach",                 "\u2713", "\u2713", "\u2713", "\u2713", "\u2713", "\u2713", "\u2713"),
]

n_rows = len(TABLE_DATA) + 1  # +header
n_cols = len(HEADERS)

tbl_shape = s.shapes.add_table(
    n_rows, n_cols,
    Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8),
)
table = tbl_shape.table

# Column widths
table.columns[0].width = Inches(2.8)
param_w = Inches((12.7 - 2.8) / (n_cols - 1))
for c in range(1, n_cols):
    table.columns[c].width = param_w

# ── Header row ──
for c, hdr in enumerate(HEADERS):
    cell = table.cell(0, c)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    run(p, hdr, BODY, 11, WHITE, bold=True)
    cell.fill.solid()
    cell.fill.fore_color.rgb = TEAL
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)

# ── Data rows ──
for r, row_data in enumerate(TABLE_DATA, start=1):
    is_ours = row_data[0] == "Our Approach"
    for c, val in enumerate(row_data):
        cell = table.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]

        if c == 0:
            p.alignment = PP_ALIGN.LEFT
            run(p, val, BODY, 11, TEAL if is_ours else BLACK, bold=is_ours)
        else:
            p.alignment = PP_ALIGN.CENTER
            if val == "\u2713":
                run(p, "\u2713", BODY, 16, CHECK_GREEN, bold=True)
            else:
                run(p, "\u2014", BODY, 14, RGBColor(0xCC, 0xCC, 0xCC))

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)

        # Row shading
        if is_ours:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREEN
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(str(OUTPUT))
print(f"Saved: {OUTPUT}")
