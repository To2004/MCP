"""Score calls_userid.csv with MISUSE risk and build a formula PPTX.

Formula (from heatmap.xlsx):
    Impact     = min(Sens * Blast * Ctx, 25)
    Likelihood = BaseDev * Anomaly * Exposure
    MISUSE     = Impact * Likelihood * Irreversibility
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOWNLOADS = Path("C:/Users/user/Downloads")
SRC_CSV = DOWNLOADS / "calls_userid.csv"
OUT_CSV = DOWNLOADS / "calls_userid_scored.csv"
OUT_XLSX = DOWNLOADS / "calls_userid_scored.xlsx"
OUT_PPTX = DOWNLOADS / "misuse_formula.pptx"

# --------------------------------------------------------------------- lookups
# These mirror the three tables on the "Lookups" sheet of heatmap.xlsx.
# Items marked (heatmap-only) appear in the Scored_Log but not the Lookups
# tables; values are taken directly from the worked rows in that sheet.

SENS_BY_EXT = {
    # ASSET TABLE — literal from Lookups
    ".sys": 5, ".exe": 5,
    ".bash": 4, ".code": 4, ".sql": 4, ".xlsx": 4, ".docx": 4,
    ".pdf": 3, ".csv": 3,
    ".md": 2, ".png": 2,
    ".txt": 1,
    # Extensions present in CSV but not in Lookups — assigned by analogy
    ".pem": 5,   # heatmap-only: private key, treated as Critical
    ".sh": 4,    # analogue of .bash
    ".c": 4,    # analogue of .code (source)
    ".json": 3,  # analogue of .csv (structured data)
    ".jpg": 2, ".jpeg": 2,  # analogue of .png
    ".xyz": 2,   # unknown binary → low-medium default
}

BLAST_BY_TOOL = {
    # TOOL TABLE — literal from Lookups
    "write_file": 5,
    "edit_file": 4,
    "move_file": 4,
    "read_file": 2,
    "list_dir": 3,
    "search": 3,
    "create_dir": 1,
    "get_file_info": 1,
    # CSV tool variants mapped onto the table families
    "read_text_file": 2,
    "read_media_file": 2,
    "read_multiple_files": 3,         # multi-file read → larger blast than single read
    "list_directory": 3,
    "list_directory_with_sizes": 3,
    "directory_tree": 3,
    "list_allowed_directories": 1,    # meta-list, no asset content
    "search_files": 3,
    "create_directory": 1,
    # Heatmap-only tools (taken from Scored_Log rows)
    "delete_file": 5,
    "chmod": 4,
    # Extras not in heatmap; conservative defaults
    "copy_file": 3,
    "execute_shell": 5,
}

# Irreversibility tiers from the Scored_Log notes (B9-B11):
#   τ1 read-only (1), τ2 modify-recoverable (2), τ3 destructive (3)
IRREV_BY_TOOL = {
    # τ3 — destructive
    "delete_file": 3,
    "execute_shell": 3,
    # τ2 — modify, generally recoverable
    "write_file": 2,
    "edit_file": 2,
    "move_file": 2,
    "chmod": 2,
    "copy_file": 2,
    "create_directory": 2,
    # τ1 — read-only (default)
}

# DIRECTORY TABLE — keys are the actual folder names under corp_filesystem_sim
# mapped to the display labels used in Lookups.
CTX_BY_DIR = {
    "security":    1.5,  # "Security Evidence"   (sub-dir of sensitive/)
    "sensitive":   1.5,  # "Sensitive Docs"
    "source_code": 1.3,  # "Source Code"
    "projects":    1.2,  # "Shared Proj Dir"
    "eval":        1.0,  # "Eval Data"
    "qa":          1.0,  # "QA Test Plans"
    "onboarding":  0.8,  # "Onboarding"
    "public":      0.6,  # "Public"
}

# Exposure is the behavioural counterpart of Ctx — how exposed the asset is
# given the access pattern. Same ordering as Ctx, normalised to 0–1.
EXPO_BY_DIR = {
    "security":    0.9,
    "sensitive":   0.9,
    "source_code": 0.8,
    "projects":    0.7,
    "eval":        0.6,
    "qa":          0.6,
    "onboarding":  0.5,
    "public":      0.3,
}

# Category-derived behavioural multipliers (calibrated to heatmap.xlsx).
# In the reference sheet Anomaly is held at 0.95 across rows -- it's an
# always-on access-pattern score. BaseDev is the one that splits dev vs
# attacker persona (0.2 vs 0.7 in the heatmap).
BASEDEV_BY_CATEGORY = {
    "VALID":     0.20,   # matches "Bob (Dev)" rows in heatmap
    "EDGE":      0.30,
    "BAD_TOOL":  0.40,
    "MISUSE":    0.60,
    "MALICIOUS": 0.70,   # matches "Mallory (Attacker)" rows in heatmap
}
ANOMALY_BY_CATEGORY = {
    "VALID":     0.95,
    "EDGE":      0.95,
    "BAD_TOOL":  0.95,
    "MISUSE":    0.95,
    "MALICIOUS": 0.95,
}
# Exposure floor so unknown-dir actions still register a baseline likelihood.
EXPO_FLOOR = 0.6


def band(score: float) -> str:
    if score >= 40:
        return "Critical"
    if score >= 20:
        return "High"
    if score >= 10:
        return "Medium"
    return "Low"


# When the path doesn't point at a specific file, infer the asset's typical
# sensitivity from the directory it lives in.
DIR_BASELINE_SENS = {
    "security":    5,
    "sensitive":   5,
    "source_code": 4,
    "projects":    3,
    "eval":        3,
    "qa":          2,
    "onboarding":  2,
    "public":      1,
}


def _norm(p: str) -> str:
    return p.lower().replace("\\\\", "/").replace("\\", "/")


def _ext_of(path: str) -> str | None:
    m = re.findall(r"\.([a-zA-Z0-9]{1,5})(?=[\"\\/]|$)", path)
    return ("." + m[-1].lower()) if m else None


def _dir_of(path: str) -> str | None:
    """Deepest matching folder key from CTX_BY_DIR inside a single path."""
    s = _norm(path)
    best_key: str | None = None
    best_pos = -1
    for key in CTX_BY_DIR:
        for m in re.finditer(rf"(?:^|[/\"\s]){re.escape(key)}(?:[/\"\s,]|$)", s):
            if m.start() > best_pos:
                best_pos = m.start()
                best_key = key
    return best_key


# Tools that carry an asset across paths (source vs destination).
MULTI_PATH_TOOLS = {"move_file", "copy_file"}


def extract_paths(args: str) -> list[str]:
    """Pull every quoted path-like value out of an args JSON blob."""
    # Match any quoted value that contains a path separator
    return re.findall(r'"([^"]*[\\\\/][^"]*)"', args)


def score_row(row: dict) -> dict:
    tool = row["tool"]
    args = row.get("args", "") or ""
    category = row.get("category", "VALID")

    ext = extract_ext(args)
    dir_key = classify_dir(args)

    if ext is not None:
        sens = SENS_BY_EXT.get(ext, 2)
        sens_reason = f"ext {ext}"
    else:
        # No file targeted (directory or meta op) — fall back to directory
        # baseline sensitivity so a list of sensitive/ still scores like
        # touching sensitive content.
        sens = DIR_BASELINE_SENS.get(dir_key, 1)
        sens_reason = f"dir baseline ({dir_key or 'no dir'})"

    blast = BLAST_BY_TOOL.get(tool, 2)
    irrev = IRREV_BY_TOOL.get(tool, 1)  # default τ1 read-only

    ctx = CTX_BY_DIR.get(dir_key, 1.0)
    expo = EXPO_BY_DIR.get(dir_key, EXPO_FLOOR)

    base_dev = BASEDEV_BY_CATEGORY.get(category, 0.3)
    anomaly = ANOMALY_BY_CATEGORY.get(category, 0.4)

    impact = min(sens * blast * ctx, 25.0)
    likelihood = base_dev * anomaly * expo
    # Floor at 1.0: any logged MCP call counts as at least baseline activity,
    # matching the heatmap reference where the minimum sampled MISUSE is 7.3
    # (their rows are all real-asset ops; ours include no-op metadata calls).
    misuse = round(max(impact * likelihood * irrev, 1.0), 2)

    why_impact = (
        f"Sens={sens} ({sens_reason}) x Blast={blast} ({tool}) "
        f"x Ctx={ctx} ({dir_key or 'unknown dir'}) -> cap 25"
    )
    why_likelihood = (
        f"BaseDev={base_dev} ({category}) x Anomaly={anomaly} "
        f"x Expo={expo} ({dir_key or 'unknown dir'}); Irrev x{irrev}"
    )

    return {
        **row,
        "sens": sens,
        "blast": blast,
        "ctx": ctx,
        "irrev": irrev,
        "base_dev": base_dev,
        "anomaly": anomaly,
        "expo": expo,
        "impact": round(impact, 2),
        "likelihood": round(likelihood, 3),
        "misuse": misuse,
        "band": band(misuse),
        "why_impact": why_impact,
        "why_likelihood": why_likelihood,
    }


# ---------------------------------------------------------------- scoring pass

def _open_for_write(path: Path):
    """Open path for writing; if locked (e.g. open in Excel), fall back to
    a sibling with a suffix and return that path instead."""
    try:
        return path, path.open("w", encoding="utf-8", newline="")
    except PermissionError:
        alt = path.with_name(path.stem + "_v2" + path.suffix)
        print(f"  {path.name} locked, writing to {alt.name} instead")
        return alt, alt.open("w", encoding="utf-8", newline="")


def write_scored_csv() -> list[dict]:
    with SRC_CSV.open(encoding="utf-8", newline="") as f:
        reader = csv.DictReader(f)
        rows = [score_row(r) for r in reader]
    if not rows:
        return []
    fieldnames = list(rows[0].keys())
    path, f = _open_for_write(OUT_CSV)
    try:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)
    finally:
        f.close()
    print(f"Wrote {path} ({len(rows)} rows)")
    return rows


def write_scored_xlsx(rows: list[dict]) -> None:
    """Write the same data to xlsx with the MISUSE column bolded."""
    import openpyxl
    from openpyxl.styles import Alignment, Font, PatternFill

    if not rows:
        return
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Scored"

    headers = list(rows[0].keys())
    misuse_col = headers.index("misuse") + 1  # 1-indexed for xlsx
    band_col = headers.index("band") + 1

    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor="2C3E50")
    misuse_header_fill = PatternFill("solid", fgColor="C03A2B")

    for j, h in enumerate(headers, start=1):
        c = ws.cell(row=1, column=j, value=h)
        c.font = header_font
        c.fill = misuse_header_fill if j == misuse_col else header_fill
        c.alignment = Alignment(horizontal="center", vertical="center")

    band_fills = {
        "Critical": PatternFill("solid", fgColor="F4C7C3"),
        "High":     PatternFill("solid", fgColor="FBE0C4"),
        "Medium":   PatternFill("solid", fgColor="FFF2BF"),
        "Low":      PatternFill("solid", fgColor="DDEBD8"),
    }

    bold = Font(bold=True)
    for i, row in enumerate(rows, start=2):
        for j, h in enumerate(headers, start=1):
            c = ws.cell(row=i, column=j, value=row[h])
            if j == misuse_col:
                c.font = bold
                c.alignment = Alignment(horizontal="center")
            if j == band_col:
                c.fill = band_fills.get(row[h], PatternFill())
                c.font = bold
                c.alignment = Alignment(horizontal="center")

    # widths
    widths = {"timestamp": 19, "args": 40, "result": 40,
              "why_impact": 55, "why_likelihood": 55}
    for j, h in enumerate(headers, start=1):
        ws.column_dimensions[ws.cell(row=1, column=j).column_letter].width = widths.get(h, 12)
    ws.freeze_panes = "A2"

    target = OUT_XLSX
    try:
        wb.save(target)
    except PermissionError:
        target = OUT_XLSX.with_name(OUT_XLSX.stem + "_v2" + OUT_XLSX.suffix)
        print(f"  {OUT_XLSX.name} locked, writing to {target.name} instead")
        wb.save(target)
    print(f"Wrote {target} ({len(rows)} rows, MISUSE column bolded)")


# --------------------------------------------------------------- pptx builders

# Palette (matches reference images)
BLUE = RGBColor(0x2E, 0x5C, 0xB8)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xD2, 0x7A, 0x1F)
RED = RGBColor(0xC0, 0x3A, 0x2B)
DARK = RGBColor(0x2C, 0x3E, 0x50)
TEXT_DARK = RGBColor(0x1F, 0x2A, 0x44)
GREY = RGBColor(0x80, 0x80, 0x80)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFA)
LIGHT_GREEN_BG = RGBColor(0xE6, 0xF3, 0xEA)
LIGHT_ORANGE_BG = RGBColor(0xFB, 0xEC, 0xD8)
LIGHT_RED_BG = RGBColor(0xFA, 0xE3, 0xDF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(shape, lines: list[tuple[str, int, bool, RGBColor]]) -> None:
    tf = shape.text_frame
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"


def rounded_box(slide, left, top, width, height, fill, line_color=None, line_width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def small_x_circle(slide, cx, cy):
    d = Inches(0.28)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d // 2, cy - d // 2, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = GREY
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "x"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = GREY
    return shape


def connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = RGBColor(0xA8, 0xB8, 0xCC)
    line.line.width = Pt(1.0)
    return line


def add_textbox(slide, left, top, width, height, lines):
    tb = slide.shapes.add_textbox(left, top, width, height)
    add_text(tb, lines)
    return tb


def lenovo_footer(slide, prs_width, prs_height, right_note: str = ""):
    # Red Lenovo block
    block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), prs_height - Inches(0.55), Inches(0.85), Inches(0.32),
    )
    block.fill.solid(); block.fill.fore_color.rgb = RED
    block.line.fill.background()
    tf = block.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Lenovo"
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

    add_textbox(
        slide, Inches(1.35), prs_height - Inches(0.52), Inches(4.5), Inches(0.3),
        [("2026 Lenovo Internal. All rights reserved.", 9, False, GREY)],
    )
    if right_note:
        add_textbox(
            slide, prs_width - Inches(5.4), prs_height - Inches(0.52),
            Inches(5.0), Inches(0.3),
            [(right_note, 9, False, GREY)],
        )


# ----------------------------------------------------------- slide 1 (formula)

def build_slide_formula(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    W, H = prs.slide_width, prs.slide_height

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(11), Inches(0.6),
                [("The Misuse Score, as a Formula Graph", 28, True, TEXT_DARK)])
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(11), Inches(0.4),
                [("Read left -> right: each factor multiplies into the one beside it, "
                  "building up to the final score.", 12, False, GREY)])

    # Left column - input factors
    in_x = Inches(0.5); in_w = Inches(2.3); in_h = Inches(0.75)
    rows = [
        ("Asset Sensitivity", "scale 1 - 5", LIGHT_BLUE_BG, BLUE),
        ("Blast Radius",       "scale 1 - 5", LIGHT_BLUE_BG, BLUE),
        ("Baseline Deviation", "0.0 - 1.0",   LIGHT_GREEN_BG, GREEN),
        ("Access Anomaly",     "0.0 - 1.0",   LIGHT_GREEN_BG, GREEN),
        ("Exposure",           "0.0 - 1.0",   LIGHT_GREEN_BG, GREEN),
        ("Reversibility tier", "1x . 2x . 3x", LIGHT_ORANGE_BG, ORANGE),
    ]
    ys = [Inches(1.5 + i * 0.95) for i in range(6)]
    for (title, sub, bg, line), y in zip(rows, ys):
        box = rounded_box(slide, in_x, y, in_w, in_h, bg, line_color=line, line_width=1.0)
        add_text(box, [
            (title, 11, True, line),
            (sub, 9, False, GREY),
        ])

    # Middle column - grouped operators
    mid_x = Inches(4.1); mid_w = Inches(3.4); mid_h = Inches(1.0)
    # IMPACT (blue)
    impact_y = (ys[0] + ys[1] + in_h) // 2 - mid_h // 2
    impact = rounded_box(slide, mid_x, impact_y, mid_w, mid_h, BLUE)
    add_text(impact, [
        ("IMPACT", 16, True, WHITE),
        ("Sensitivity x Blast  ->  max 25", 10, False, WHITE),
    ])

    # LIKELIHOOD (green) - spans 3 inputs
    lik_y = (ys[2] + ys[4] + in_h) // 2 - mid_h // 2
    lik = rounded_box(slide, mid_x, lik_y, mid_w, mid_h, GREEN)
    add_text(lik, [
        ("LIKELIHOOD", 16, True, WHITE),
        ("Dev x Anomaly x Exposure", 10, False, WHITE),
    ])

    # IRREVERSIBILITY (orange)
    irr_y = ys[5] + in_h // 2 - mid_h // 2
    irr = rounded_box(slide, mid_x, irr_y, mid_w, mid_h, ORANGE)
    add_text(irr, [
        ("IRREVERSIBILITY", 14, True, WHITE),
        ("multiplier x1 -> x3", 10, False, WHITE),
    ])

    # x circles between inputs and middle
    cx_join = Inches(3.7)
    small_x_circle(slide, cx_join, (ys[0] + ys[1] + in_h) // 2)
    small_x_circle(slide, cx_join, (ys[2] + ys[3] + in_h) // 2)
    small_x_circle(slide, cx_join, (ys[3] + ys[4] + in_h) // 2)
    # connectors input -> circle -> impact/lik/irr
    for y in ys[:2]:
        connector(slide, in_x + in_w, y + in_h // 2, mid_x, impact_y + mid_h // 2)
    for y in ys[2:5]:
        connector(slide, in_x + in_w, y + in_h // 2, mid_x, lik_y + mid_h // 2)
    connector(slide, in_x + in_w, ys[5] + in_h // 2, mid_x, irr_y + mid_h // 2)

    # Right - MISUSE SCORE
    out_x = Inches(9.3); out_w = Inches(3.3); out_h = Inches(2.4)
    out_y = (impact_y + irr_y + mid_h) // 2 - out_h // 2
    out = rounded_box(slide, out_x, out_y, out_w, out_h, DARK)
    add_text(out, [
        ("MISUSE", 22, True, WHITE),
        ("SCORE", 22, True, WHITE),
        ("", 6, False, WHITE),
        ("Impact x Likelihood", 11, False, WHITE),
        ("x Irreversibility", 11, False, WHITE),
    ])

    # x circles between middle and right
    cx2 = Inches(8.6)
    small_x_circle(slide, cx2, impact_y + mid_h // 2)
    small_x_circle(slide, cx2, lik_y + mid_h // 2)
    small_x_circle(slide, cx2, irr_y + mid_h // 2)
    for y_mid in (impact_y, lik_y, irr_y):
        connector(slide, mid_x + mid_w, y_mid + mid_h // 2, out_x, out_y + out_h // 2)

    lenovo_footer(
        slide, W, H,
        right_note="Scales adapted from: NIST SP 800-30 Rev.1  -  OWASP Risk Rating Methodology  -  NIST SP 800-60 / FIPS 199",
    )


# ------------------------------------------------------------ slide 2 (example)

def build_slide_example(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(11), Inches(0.6),
                [("Reading the Graph Through One Example", 28, True, TEXT_DARK)])

    # Scenario banner
    banner = rounded_box(slide, Inches(0.5), Inches(0.95), Inches(7.0), Inches(0.45),
                        WHITE, line_color=RED, line_width=1.0)
    tf = banner.text_frame
    tf.margin_left = Emu(100000); tf.margin_top = Emu(20000)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Scenario:  "
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = RED
    r = p.add_run(); r.text = "a user sends db:delete on the financial-records table"
    r.font.size = Pt(12); r.font.color.rgb = TEXT_DARK

    # Inputs with concrete values
    in_x = Inches(0.5); in_w = Inches(2.3); in_h = Inches(0.75)
    rows = [
        ("Asset Sensitivity", "5  -  financial records", LIGHT_BLUE_BG, BLUE, True),
        ("Blast Radius",      "4  -  bulk table",        LIGHT_BLUE_BG, BLUE, True),
        ("Baseline Deviation","0.8  -  new op + asset",  WHITE, GREY,  False),
        ("Access Anomaly",    "0.7  -  off-pattern",     WHITE, GREY,  False),
        ("Exposure",          "0.9  -  sensitive DB",    WHITE, GREY,  False),
        ("Reversibility",     "3x  -  hard delete",      LIGHT_ORANGE_BG, ORANGE, True),
    ]
    ys = [Inches(1.6 + i * 0.85) for i in range(6)]
    for (title, sub, bg, line, boxed), y in zip(rows, ys):
        if boxed:
            box = rounded_box(slide, in_x, y, in_w, in_h, bg, line_color=line, line_width=1.0)
            add_text(box, [(title, 11, True, line), (sub, 9, False, GREY)])
        else:
            tb = add_textbox(slide, in_x + Inches(0.05), y, in_w, in_h,
                             [(title, 11, True, TEXT_DARK), (sub, 9, False, GREY)])

    # Middle: IMPACT, LIKELIHOOD, IRREVERSIBILITY with concrete values
    mid_x = Inches(4.1); mid_w = Inches(3.4); mid_h = Inches(0.95)

    impact_y = (ys[0] + ys[1] + in_h) // 2 - mid_h // 2
    impact = rounded_box(slide, mid_x, impact_y, mid_w, mid_h, BLUE)
    add_text(impact, [("IMPACT = 20", 18, True, WHITE), ("5 x 4", 10, False, WHITE)])

    lik_y = (ys[2] + ys[4] + in_h) // 2 - mid_h // 2
    lik = rounded_box(slide, mid_x, lik_y, mid_w, mid_h, RGBColor(0xC2, 0xCC, 0xD9))
    add_text(lik, [
        ("LIKELIHOOD = 0.50", 18, True, WHITE),
        ("0.8 x 0.7 x 0.9", 10, False, WHITE),
    ])

    irr_y = ys[5] + in_h // 2 - mid_h // 2
    irr = rounded_box(slide, mid_x, irr_y, mid_w, mid_h, ORANGE)
    add_text(irr, [("IRREVERSIBILITY = 3x", 16, True, WHITE),
                   ("no rollback", 10, False, WHITE)])

    # x circles + connectors
    cx_join = Inches(3.7)
    small_x_circle(slide, cx_join, (ys[0] + ys[1] + in_h) // 2)
    small_x_circle(slide, cx_join, (ys[2] + ys[3] + in_h) // 2)
    small_x_circle(slide, cx_join, (ys[3] + ys[4] + in_h) // 2)
    for y in ys[:2]:
        connector(slide, in_x + in_w, y + in_h // 2, mid_x, impact_y + mid_h // 2)
    for y in ys[2:5]:
        connector(slide, in_x + in_w, y + in_h // 2, mid_x, lik_y + mid_h // 2)
    connector(slide, in_x + in_w, ys[5] + in_h // 2, mid_x, irr_y + mid_h // 2)

    # MISUSE red box
    out_x = Inches(9.3); out_w = Inches(3.3); out_h = Inches(1.9)
    out_y = (impact_y + irr_y + mid_h) // 2 - out_h // 2
    out = rounded_box(slide, out_x, out_y, out_w, out_h, RED)
    add_text(out, [
        ("MISUSE = 30.0", 20, True, WHITE),
        ("20 x 0.50 x 3", 10, False, WHITE),
        ("", 6, False, WHITE),
    ])
    crit = rounded_box(slide, out_x + Inches(0.55), out_y + Inches(1.25),
                       Inches(2.2), Inches(0.45), WHITE)
    tfc = crit.text_frame; tfc.margin_top = Emu(20000)
    pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
    rc = pc.add_run(); rc.text = "CRITICAL"
    rc.font.size = Pt(14); rc.font.bold = True; rc.font.color.rgb = RED

    cx2 = Inches(8.6)
    small_x_circle(slide, cx2, impact_y + mid_h // 2)
    small_x_circle(slide, cx2, lik_y + mid_h // 2)
    small_x_circle(slide, cx2, irr_y + mid_h // 2)
    for y_mid in (impact_y, lik_y, irr_y):
        connector(slide, mid_x + mid_w, y_mid + mid_h // 2, out_x, out_y + out_h // 2)

    # Right-side callouts
    add_textbox(slide, Inches(7.7), Inches(1.55), Inches(2.4), Inches(1.2), [
        ("Why Impact is high", 10, True, TEXT_DARK),
        ("Financial records are a top-tier asset (5), and a bulk "
         "table means wide reach (4).", 9, False, GREY),
    ])
    add_textbox(slide, Inches(7.7), Inches(5.5), Inches(2.4), Inches(1.2), [
        ("The deciding factor", 10, True, ORANGE),
        ("A hard delete (3x) can't be undone. That multiplier is what "
         "tips the score into Critical.", 9, False, GREY),
    ])

    lenovo_footer(
        slide, W, H,
        right_note="Likelihood (green) is only moderate here - but a top asset destroyed irreversibly still scores Critical.",
    )


# ----------------------------------------------------------------------- main

def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_slide_formula(prs)
    build_slide_example(prs)
    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    scored = write_scored_csv()
    write_scored_xlsx(scored)
    build_pptx()
