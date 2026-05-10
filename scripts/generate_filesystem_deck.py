"""Generate the MCP Filesystem Tools teaching deck (18 slides)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
BG = RGBColor(0x1E, 0x1E, 0x2E)       # dark navy background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x89, 0xB4, 0xFA)   # blue accent
GREEN = RGBColor(0xA6, 0xE3, 0xA1)    # good example
RED = RGBColor(0xF3, 0x8B, 0xA8)      # bad example / error
YELLOW = RGBColor(0xF9, 0xE2, 0xAF)   # warning / edge case
GRAY = RGBColor(0x45, 0x47, 0x5A)     # subtle box fill
MONO_FONT = "Courier New"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUTPUT = Path(__file__).parent.parent / "presentations" / "mcp-filesystem-tools.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs: Presentation) -> Slide:
    """Add a blank slide and paint the background."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def text_box(
    slide: Slide,
    left,
    top,
    width,
    height,
    text,
    size=Pt(18),
    bold=False,
    color=WHITE,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    wrap=True,
):
    """Add a text box and return the shape."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return txb


def code_box(slide: Slide, left, top, width, height, code, fill=GRAY, text_color=WHITE):
    """A monospace code block with a filled background rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code
    run.font.size = Pt(12)
    run.font.name = MONO_FONT
    run.font.color.rgb = text_color
    return shape


def label_box(
    slide: Slide,
    left,
    top,
    width,
    height,
    text,
    fill=ACCENT,
    text_color=BG,
    size=Pt(14),
):
    """Colored label / badge box."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.name = BODY_FONT
    run.font.color.rgb = text_color
    return shape


def arrow(slide: Slide, x1, y1, x2, y2, color=ACCENT):
    """Draw a straight connector between two points."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        x1, y1, x2, y2,
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def slide_title(slide: Slide, title, subtitle=None) -> None:
    """Add a slide title (large) and optional subtitle."""
    text_box(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
             title, size=Pt(32), bold=True, color=ACCENT)
    if subtitle:
        text_box(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
                 subtitle, size=Pt(18), color=WHITE)


# ---------------------------------------------------------------------------
# Slide stubs — replaced in Tasks 3, 4, 5
# ---------------------------------------------------------------------------

def _slide_01_title(prs: Presentation) -> None:
    slide = add_slide(prs)
    text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
             "MCP Filesystem Server", size=Pt(48), bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    text_box(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(0.8),
             "How It Works — Tool Reference & Teaching Guide",
             size=Pt(24), color=WHITE, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
             "@modelcontextprotocol/server-filesystem  v1.27.0",
             size=Pt(16), color=GRAY, align=PP_ALIGN.CENTER)


def _slide_02_what_is_mcp(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "What is MCP?",
                "Model Context Protocol — lets AI agents call tools on external servers")

    boxes = [
        (Inches(0.5),  "AI Agent",    ACCENT),
        (Inches(3.5),  "JSON-RPC 2.0", YELLOW),
        (Inches(6.5),  "MCP Server",   GREEN),
        (Inches(9.5),  "Filesystem",   GRAY),
    ]
    for left, label, color in boxes:
        label_box(slide, left, Inches(3.2), Inches(2.5), Inches(1.0),
                  label, fill=color, text_color=BG, size=Pt(16))

    for i in range(3):
        x1 = Inches(0.5 + i * 3 + 2.5)
        x2 = Inches(0.5 + (i + 1) * 3)
        arrow(slide, x1, Inches(3.7), x2, Inches(3.7))

    text_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.2),
             "An agent sends a JSON-RPC request naming a tool (e.g. read_text_file) with "
             "parameters. The MCP server executes it and returns a result. "
             "The filesystem server exposes 13 tools — all sandboxed to an allowed root.",
             size=Pt(16), color=WHITE)


def _slide_03_request_response(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "Request → Response Flow")

    request_json = (
        '{\n'
        '  "jsonrpc": "2.0",\n'
        '  "id": 1,\n'
        '  "method": "tools/call",\n'
        '  "params": {\n'
        '    "name": "read_text_file",\n'
        '    "arguments": {\n'
        '      "path": "sensitive/security/audit_log.txt"\n'
        '    }\n'
        '  }\n'
        '}'
    )
    response_json = (
        '{\n'
        '  "jsonrpc": "2.0",\n'
        '  "id": 1,\n'
        '  "result": {\n'
        '    "content": [{\n'
        '      "type": "text",\n'
        '      "text": "2026-01-15 09:00 | LOGIN..."\n'
        '    }]\n'
        '  }\n'
        '}'
    )

    text_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.4),
             "REQUEST", size=Pt(14), bold=True, color=GREEN)
    code_box(slide, Inches(0.5), Inches(1.9), Inches(5.5), Inches(4.8), request_json)

    arrow(slide, Inches(6.2), Inches(4.2), Inches(7.0), Inches(4.2))

    text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
             "RESPONSE", size=Pt(14), bold=True, color=ACCENT)
    code_box(slide, Inches(7.2), Inches(1.9), Inches(5.5), Inches(4.8), response_json)


def _slide_04_security_boundary(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "The Security Boundary",
                "Every path must live inside the allowed root — no exceptions")

    root_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.5))
    root_box.fill.solid()
    root_box.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x1A)
    root_box.line.color.rgb = GREEN
    root_box.line.width = Pt(2)

    text_box(slide, Inches(0.7), Inches(1.9), Inches(5.5), Inches(0.4),
             "✓  Allowed root: /corp_filesystem/", size=Pt(14), bold=True, color=GREEN)

    allowed_paths = [
        "sensitive/security/audit_log.txt",
        "projects/known_defects.csv",
        "source_code/core.c",
    ]
    for i, path in enumerate(allowed_paths):
        text_box(slide, Inches(0.9), Inches(2.5 + i * 0.5), Inches(5.5), Inches(0.45),
                 f"  {path}", size=Pt(13), color=WHITE)

    blocked = [
        ("../etc/passwd",     "Access denied - path outside allowed directories"),
        ("C:/Windows/win.ini","Access denied - path outside allowed directories"),
        ("(empty string)",    "Access denied - path outside allowed directories"),
    ]
    for i, (attempt, error) in enumerate(blocked):
        y = Inches(2.1 + i * 1.1)
        label_box(slide, Inches(7.0), y, Inches(5.8), Inches(0.45),
                  f"✗  {attempt}", fill=RED, text_color=WHITE, size=Pt(12))
        text_box(slide, Inches(7.0), y + Inches(0.5), Inches(5.8), Inches(0.4),
                 f"→ {error}", size=Pt(11), color=RED)


def _slide_05_read_overview(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "Read Tools")

    tools = [
        ("read_text_file",      "Returns file content as a UTF-8 string.\nBest for: .txt .md .csv .sql .c .sh .pem"),
        ("read_media_file",     "Returns binary as base64 blob + MIME type.\nBest for: .png .pdf .docx .xlsx .exe .sys"),
        ("read_multiple_files", "Reads an array of files in one call.\nPartial failure doesn't abort the batch."),
    ]
    for i, (name, desc) in enumerate(tools):
        x = Inches(0.5 + i * 4.2)
        label_box(slide, x, Inches(2.0), Inches(3.9), Inches(0.55),
                  name, fill=ACCENT, text_color=BG, size=Pt(14))
        text_box(slide, x, Inches(2.7), Inches(3.9), Inches(1.5),
                 desc, size=Pt(14), color=WHITE)


def _slide_06_read_text_file(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "read_text_file",
                "Input: path (str), optional head/tail (int) → Output: plain UTF-8 string")

    text_box(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.4),
             "Works with:", size=Pt(14), bold=True, color=GREEN)
    text_box(slide, Inches(0.5), Inches(2.2), Inches(6.0), Inches(0.5),
             ".txt   .md   .csv   .sql   .c   .py   .sh   .bash   .pem",
             size=Pt(14), color=GREEN, font=MONO_FONT)

    text_box(slide, Inches(0.5), Inches(2.9), Inches(6.0), Inches(0.4),
             "Do NOT use with (binary — use read_media_file):", size=Pt(14), bold=True, color=RED)
    text_box(slide, Inches(0.5), Inches(3.3), Inches(6.0), Inches(0.5),
             ".png   .pdf   .docx   .xlsx   .exe   .sys",
             size=Pt(14), color=RED, font=MONO_FONT)

    text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.4),
             "✓ Good", size=Pt(13), bold=True, color=GREEN)
    code_box(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(1.6),
             'read_text_file(\n  path="sensitive/security/audit_log.txt"\n)\n→ "2026-01-15 09:00 | LOGIN | ..."',
             fill=RGBColor(0x1A, 0x2A, 0x1A), text_color=GREEN)

    text_box(slide, Inches(7.0), Inches(4.0), Inches(5.8), Inches(0.4),
             "✗ Bad — binary file", size=Pt(13), bold=True, color=RED)
    code_box(slide, Inches(7.0), Inches(4.4), Inches(5.8), Inches(1.6),
             'read_text_file(\n  path="onboarding/org_chart.png"\n)\n→ Error: Cannot decode as UTF-8',
             fill=RGBColor(0x2A, 0x1A, 0x1A), text_color=RED)

    text_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
             "Edge case: head=5 → first 5 lines only. tail=3 → last 3 lines. Cannot combine both.",
             size=Pt(13), color=YELLOW)


def _slide_07_read_media_file(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "read_media_file",
                "Input: path (str) → Output: base64 blob + mimeType")

    text_box(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.4),
             "Use for binary files:", size=Pt(14), bold=True, color=ACCENT)
    rows = [
        (".png / .jpg", "image", "image/png"),
        (".pdf",        "blob",  "application/pdf"),
        (".docx",       "blob",  "application/vnd.openxmlformats..."),
        (".xlsx",       "blob",  "application/vnd.openxmlformats..."),
        (".exe / .sys", "blob",  "application/octet-stream"),
    ]
    for i, (ext, typ, mime) in enumerate(rows):
        y = Inches(2.3 + i * 0.55)
        text_box(slide, Inches(0.5), y, Inches(2.0), Inches(0.5),
                 ext, size=Pt(13), color=WHITE, font=MONO_FONT)
        text_box(slide, Inches(2.6), y, Inches(1.2), Inches(0.5),
                 typ, size=Pt(13), color=ACCENT, font=MONO_FONT)
        text_box(slide, Inches(4.0), y, Inches(3.0), Inches(0.5),
                 mime, size=Pt(11), color=GRAY, font=MONO_FONT)

    text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.4),
             "✓ Good", size=Pt(13), bold=True, color=GREEN)
    code_box(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.2),
             'read_media_file(path="public/logo.png")\n\n→ {\n    "type": "image",\n    "data": "iVBORw0KGgo...",\n    "mimeType": "image/png"\n  }',
             fill=RGBColor(0x1A, 0x2A, 0x1A), text_color=GREEN)

    text_box(slide, Inches(7.0), Inches(4.6), Inches(5.8), Inches(0.4),
             "✗ Bad — text file wastes tokens", size=Pt(13), bold=True, color=YELLOW)
    code_box(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.5),
             'read_media_file(path="audit_log.txt")\n# works but returns base64 of plain text\n# use read_text_file instead',
             fill=RGBColor(0x2A, 0x28, 0x1A), text_color=YELLOW)


def _slide_08_read_multiple_files(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "read_multiple_files",
                "Input: paths[] → Output: all contents in one string (partial failure OK)")

    code_box(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(2.5),
             'read_multiple_files(paths=[\n  "sensitive/financials/payslips_q1.csv",\n  "does_not_exist.txt",\n  "projects/known_defects.csv"\n])',
             fill=GRAY)

    arrow(slide, Inches(6.7), Inches(3.0), Inches(7.3), Inches(3.0))

    code_box(slide, Inches(7.5), Inches(1.8), Inches(5.3), Inches(2.5),
             'payslips_q1.csv:\nemployee_id,name,salary\n101,Alice,72000\n\n---\ndoes_not_exist.txt: Error - ENOENT\n\n---\nknown_defects.csv:\nid,title,severity\n1,Memory leak,high',
             fill=GRAY)

    label_box(slide, Inches(7.5), Inches(4.5), Inches(5.3), Inches(0.45),
              "Bad path → inline error, rest succeed", fill=YELLOW, text_color=BG, size=Pt(13))

    text_box(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(1.4),
             "Key points:\n"
             "• Response is one string — parse it yourself\n"
             "• No atomicity — some succeed, some fail\n"
             "• Check for 'Error -' lines to detect failures",
             size=Pt(14), color=WHITE)


def _slide_09_write_tools(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "Write Tools")

    label_box(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.55),
              "write_file", fill=RED, text_color=WHITE)
    text_box(slide, Inches(0.5), Inches(2.5), Inches(5.5), Inches(0.4),
             "Creates or OVERWRITES a file — no confirmation", size=Pt(13), color=WHITE)
    code_box(slide, Inches(0.5), Inches(3.0), Inches(5.5), Inches(1.3),
             'write_file(\n  path="source_code/notes.txt",\n  content="Sprint 12 goals:\\n- Fix auth"\n)\n→ "Successfully wrote to notes.txt"',
             fill=GRAY)
    label_box(slide, Inches(0.5), Inches(4.5), Inches(5.5), Inches(0.4),
              "⚠ destructive=true  idempotent=true", fill=YELLOW, text_color=BG, size=Pt(12))

    label_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.55),
              "edit_file", fill=ACCENT, text_color=BG)
    text_box(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(0.4),
             "Find-and-replace edits → returns git diff", size=Pt(13), color=WHITE)
    code_box(slide, Inches(7.0), Inches(3.0), Inches(5.8), Inches(2.5),
             'edit_file(\n  path="notes.txt",\n  edits=[{\n    "oldText": "Fix auth",\n    "newText": "Fix auth module"\n  }]\n)\n→ "@@ -1 +1 @@\\n-Fix auth\\n+Fix auth module"',
             fill=GRAY)
    label_box(slide, Inches(7.0), Inches(5.6), Inches(5.8), Inches(0.4),
              "⚠ destructive=true  idempotent=false", fill=YELLOW, text_color=BG, size=Pt(12))

    text_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.45),
             "Tip: use dryRun=true to preview the diff before committing changes",
             size=Pt(13), color=YELLOW)


def _slide_10_list_tools(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "Navigation — List & Tree Tools")

    cols = [
        ("list_directory", ACCENT,
         "sensitive/\n[DIR] contracts\n[DIR] financials\n[DIR] security"),
        ("list_directory_with_sizes", GREEN,
         "sensitive/financials/\n[FILE] payslips_q1.csv  74 B\n[FILE] budget_2026.xlsx  15 B\n\nTotal: 2 files, 0 dirs\nCombined size: 89 B"),
        ("directory_tree", YELLOW,
         '[\n  {"name": "contracts",\n   "type": "directory",\n   "children": [...]},\n  {"name": "financials",\n   "type": "directory",\n   "children": [...]}\n]'),
    ]
    for i, (name, color, output) in enumerate(cols):
        x = Inches(0.3 + i * 4.35)
        label_box(slide, x, Inches(1.8), Inches(4.0), Inches(0.5),
                  name, fill=color, text_color=BG, size=Pt(13))
        code_box(slide, x, Inches(2.4), Inches(4.0), Inches(3.8),
                 output, fill=GRAY)

    text_box(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.45),
             "directory_tree output is a JSON array string — parse it. Use excludePatterns=[\"*.log\"] to skip noisy files.",
             size=Pt(13), color=YELLOW)


def _slide_11_search_files(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "search_files — Glob Pattern Cheatsheet")

    rows = [
        ("*",             "Current dir only",    "NO",  'README.md  (only if in root)\n→ returns "No matches found" if none'),
        ("**/*",          "All files recursive", "YES", "sensitive/financials/payslips_q1.csv\nprojects/known_defects.csv\n..."),
        ("**/*.csv",      "CSV files only",      "YES", "sensitive/financials/payslips_q1.csv\nprojects/known_defects.csv"),
        ("sensitive/**",  "Inside sensitive/",   "YES", "sensitive/contracts/...\nsensitive/financials/..."),
    ]
    headers = ["Pattern", "Meaning", "Recursive?", "Example result"]
    col_x = [Inches(0.3), Inches(2.8), Inches(5.5), Inches(7.2)]
    col_w = [Inches(2.3), Inches(2.5), Inches(1.5), Inches(5.7)]

    for j, h in enumerate(headers):
        text_box(slide, col_x[j], Inches(1.8), col_w[j], Inches(0.45),
                 h, size=Pt(13), bold=True, color=ACCENT)

    for i, (pat, meaning, rec, example) in enumerate(rows):
        y = Inches(2.4 + i * 1.1)
        bg = RGBColor(0x1A, 0x1A, 0x3A) if i % 2 == 0 else GRAY
        vals = [pat, meaning, rec, example]
        colors = [RED if pat == "*" else GREEN, WHITE, GREEN if rec == "YES" else YELLOW, WHITE]
        for j, (val, col) in enumerate(zip(vals, colors)):
            code_box(slide, col_x[j], y, col_w[j], Inches(1.0),
                     val, fill=bg, text_color=col)

    text_box(slide, Inches(0.3), Inches(6.8), Inches(12.5), Inches(0.45),
             "⚠  * is NOT recursive — the #1 beginner mistake. Always use **/* for recursive search.",
             size=Pt(14), bold=True, color=RED)


def _slide_12_info_admin(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "Info & Admin Tools")

    tools = [
        ("get_file_info",
         "path → metadata\n(size, timestamps,\ntype, permissions)",
         "size: 134\nmodified: 2026-01-15\nisFile: true\npermissions: 644"),
        ("create_directory",
         "path → creates dir\n(and parents).\nSafe to retry.",
         'create_directory(\n  "source_code/feature/utils"\n)\n→ "Successfully created..."'),
        ("move_file",
         "source, destination\n→ moves or renames.\nBehavior on existing\ndest varies by OS.",
         'move_file(\n  source="notes.txt",\n  destination="archive/notes.txt"\n)\n→ "Successfully moved..."'),
        ("list_allowed_directories",
         "No params → lists\nserver root dirs.\nCall this first.",
         '→ "Allowed directories:\n/corp_filesystem_sim"'),
    ]
    for i, (name, desc, example) in enumerate(tools):
        x = Inches(0.2 + i * 3.3)
        label_box(slide, x, Inches(1.8), Inches(3.0), Inches(0.5),
                  name, fill=ACCENT, text_color=BG, size=Pt(12))
        text_box(slide, x, Inches(2.4), Inches(3.0), Inches(1.3),
                 desc, size=Pt(12), color=WHITE)
        code_box(slide, x, Inches(3.8), Inches(3.0), Inches(2.5),
                 example, fill=GRAY)


def _slide_13_file_type_matrix(prs: Presentation) -> None:
    slide = add_slide(prs)
    slide_title(slide, "File Type × Tool Matrix")

    rows = [
        (".txt / .md / .csv", "✅", "⚠️ wasteful", "✅", "✅"),
        (".sql / .c / .sh",   "✅", "⚠️ wasteful", "✅", "✅"),
        (".pem",              "✅", "⚠️ wasteful", "✅", "✅"),
        (".png",              "❌",  "✅",          "❌",  "✅"),
        (".pdf",              "❌",  "✅",          "❌",  "✅"),
        (".docx / .xlsx",     "❌",  "✅",          "❌",  "✅"),
        (".exe / .sys",       "❌",  "✅ (opaque)", "❌",  "✅"),
    ]
    cols = ["File type", "read_text_file", "read_media_file", "write_file", "search_files"]
    col_x = [Inches(0.2), Inches(2.8), Inches(5.1), Inches(8.2), Inches(10.5)]
    col_w = [Inches(2.4), Inches(2.1), Inches(2.9), Inches(2.1), Inches(2.5)]

    for j, h in enumerate(cols):
        label_box(slide, col_x[j], Inches(1.8), col_w[j], Inches(0.45),
                  h, fill=GRAY, text_color=ACCENT, size=Pt(12))

    for i, (ftype, *cells) in enumerate(rows):
        y = Inches(2.4 + i * 0.6)
        text_box(slide, col_x[0], y, col_w[0], Inches(0.55),
                 ftype, size=Pt(12), color=WHITE, font=MONO_FONT)
        for j, cell in enumerate(cells):
            color = GREEN if "✅" in cell else (RED if "❌" in cell else YELLOW)
            text_box(slide, col_x[j + 1], y, col_w[j + 1], Inches(0.55),
                     cell, size=Pt(13), color=color, align=PP_ALIGN.CENTER)

    text_box(slide, Inches(0.2), Inches(6.8), Inches(12.5), Inches(0.45),
             "✅ correct tool  ·  ⚠️ works but suboptimal  ·  ❌ wrong tool / fails",
             size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)
def _slide_14_glob_patterns(prs): add_slide(prs)
def _slide_15_path_pitfalls(prs): add_slide(prs)
def _slide_16_silent_overwrite(prs): add_slide(prs)
def _slide_17_partial_batch(prs): add_slide(prs)
def _slide_18_missing_tools(prs): add_slide(prs)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def generate_deck() -> Presentation:
    prs = new_presentation()
    _slide_01_title(prs)
    _slide_02_what_is_mcp(prs)
    _slide_03_request_response(prs)
    _slide_04_security_boundary(prs)
    _slide_05_read_overview(prs)
    _slide_06_read_text_file(prs)
    _slide_07_read_media_file(prs)
    _slide_08_read_multiple_files(prs)
    _slide_09_write_tools(prs)
    _slide_10_list_tools(prs)
    _slide_11_search_files(prs)
    _slide_12_info_admin(prs)
    _slide_13_file_type_matrix(prs)
    _slide_14_glob_patterns(prs)
    _slide_15_path_pitfalls(prs)
    _slide_16_silent_overwrite(prs)
    _slide_17_partial_batch(prs)
    _slide_18_missing_tools(prs)
    return prs


if __name__ == "__main__":
    prs = generate_deck()
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
