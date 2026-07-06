"""Generate a 3-slide deck summarizing the MCP asset scanner."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT = Path.home() / "Downloads" / "mcp_scanner_overview.pptx"

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x86, 0x4E)


def add_title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_text_box(slide, left, top, width, height, lines, size=16, bold_first=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = GREY
        if bold_first and i == 0:
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.font.size = Pt(size + 2)
    return tb


def add_pill(slide, left, top, width, height, text, fill, font_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill
    tf = s.text_frame
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 2  # center
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = font_color or RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_code_block(slide, left, top, width, height, lines):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0xDC, 0xDC, 0xDC)


def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "MCP Asset Scanner"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tb2 = s.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(1))
    p = tb2.text_frame.paragraphs[0]
    p.text = "Enumerate and rank what every connected MCP server can reach"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCF, 0xDC, 0xE8)

    bullets = [
        "Connects to each server in ~/.claude.json (read-only)",
        "Generic protocol enumeration + per-kind fast paths (filesystem, sqlite)",
        "Deterministic discovery -> agentic understanding (local LLM, Qwen2.5)",
        "Output: Rank | Name | Risk Level | Reasoning table",
    ]
    box = s.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(12), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.space_after = Pt(8)


def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "How it works — pipeline")

    # Pipeline pills
    y = Inches(1.5)
    h = Inches(0.7)
    w = Inches(2.4)
    gap = Inches(0.3)
    labels = ["config_reader", "enumerate", "rank", "table"]
    colors = [ACCENT, ACCENT, ACCENT, GREEN]
    x = Inches(0.6)
    for label, c in zip(labels, colors):
        add_pill(s, x, y, w, h, label, c)
        x += w + gap

    # Arrows row (simple dashes via textbox)
    arrows = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12), Inches(0.7))
    p = arrows.text_frame.paragraphs[0]
    p.text = "                    →                    →                    →"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Two columns: discovery vs understanding
    col_top = Inches(2.7)
    col_h = Inches(4.3)
    col_w = Inches(6.0)

    left_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.6), col_top, col_w, col_h)
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT
    left_box.line.color.rgb = ACCENT
    tf = left_box.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Discovery (deterministic)"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = NAVY
    for line in [
        "",
        "• Generic: MCP list_resources + read-only tools",
        "• Filesystem: os.walk, directory_tree, list_directory",
        "• SQLite: list_tables + describe_table",
        "• Unreachable servers: resolver (web/GitHub → LLM,",
        "   else theorise plausible assets)",
        "• safety.py gate: only LIST / SEARCH / METADATA",
        "   tools may be invoked. Never write/delete.",
    ]:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(15)
        pp.font.color.rgb = GREY

    right_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(6.9), col_top, col_w, col_h)
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT
    right_box.line.color.rgb = GREEN
    tf = right_box.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Understanding (agentic)"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = NAVY
    for line in [
        "",
        "• Local LLM (Qwen2.5:32b via Ollama, no cloud)",
        "• Reads server's own tool/resource descriptions",
        "   (AssetInventory.context)",
        "• Ranks against FIPS / MITRE / OWASP sensitivity",
        "   anchors → Risk Level + Reasoning",
        "• Unknown server kinds work with zero per-kind code",
        "• If Ollama is down: anchored rows still print,",
        "   rest flagged 'unranked'",
    ]:
        pp = tf.add_paragraph()
        pp.text = line
        pp.font.size = Pt(15)
        pp.font.color.rgb = GREY


def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Examples")

    # CLI block
    add_text_box(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
                 ["CLI"], size=18, bold_first=True)
    add_code_block(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.7), [
        "$ python -m mcp_security.scanner --root demo/corp_filesystem",
        "$ python -m mcp_security.scanner --root demo/corp_sqlite/corp.db --kind sqlite",
        "$ python -m mcp_security.scanner               # every server in ~/.claude.json",
        "$ python -m mcp_security.scanner --server github --out report.md --no-web",
    ])

    # Sample output header
    add_text_box(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.4),
                 ["Sample ranked output  (Rank | Name | Risk Level | Reasoning)"],
                 size=18, bold_first=True)

    # Table
    rows = [
        ("Rank", "Name", "Risk Level", "Reasoning"),
        ("1", "corp_filesystem / hr/salaries/", "CRITICAL",
         "PII + compensation; FIPS-199 High confidentiality"),
        ("2", "corp_sqlite / customers table", "HIGH",
         "Customer records w/ emails; OWASP A02 exposure"),
        ("3", "github / private repo: infra-terraform", "HIGH",
         "Cloud credentials + IaC; MITRE T1552"),
        ("4", "slack / #eng-incidents channel", "MEDIUM",
         "Operational secrets in chat history"),
        ("5", "corp_filesystem / public/marketing/", "LOW",
         "Marketing assets, intended public"),
    ]
    table_left = Inches(0.6)
    table_top = Inches(4.0)
    table_w = Inches(12.1)
    table_h = Inches(3.0)
    table = s.shapes.add_table(rows=len(rows), cols=4,
                               left=table_left, top=table_top,
                               width=table_w, height=table_h).table
    col_widths = [Inches(0.6), Inches(3.6), Inches(1.6), Inches(6.3)]
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    risk_color = {
        "CRITICAL": RED, "HIGH": RGBColor(0xE6, 0x7E, 0x22),
        "MEDIUM": RGBColor(0xD4, 0xAC, 0x0D), "LOW": GREEN,
    }
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(13)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    elif c == 2 and val in risk_color:
                        run.font.bold = True
                        run.font.color.rgb = risk_color[val]
                    else:
                        run.font.color.rgb = GREY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide1(prs)
    slide2(prs)
    slide3(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
